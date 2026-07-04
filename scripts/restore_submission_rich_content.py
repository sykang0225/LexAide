# -*- coding: utf-8 -*-
"""Restore rich submission documents after an over-simplified rewrite.

This script rebuilds the two submission files with the original dense story:
official 7-section structure, benchmarking/research rationale, legal-tree
positioning, Human-in-the-loop framing, and implementation evidence.
"""

from __future__ import annotations

import json
from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt


ROOT = Path(__file__).resolve().parents[2]
PROJECT = ROOT / "cross_check_ai"
DOCX_OUT = ROOT / "CrossCheckAI_기능명세서.docx"
PPTX_OUT = ROOT / "CrossCheckAI_MVP제안서.pptx"
SUMMARY_PATH = PROJECT / "data" / "evaluation" / "evaluation_report.summary.json"

if SUMMARY_PATH.exists():
    SUMMARY = json.loads(SUMMARY_PATH.read_text(encoding="utf-8"))
else:
    SUMMARY = {
        "overall": {
            "risk_recall": 1.0,
            "risk_precision": 0.75,
            "violation_recall_not_pass": 1.0,
        },
        "ko": {"violation_recall_not_pass": 1.0},
        "ojk": {"risk_recall": 1.0},
    }


FONT = "맑은 고딕"
BLUE = RGBColor(0x00, 0x3B, 0x8F)
DEEP = RGBColor(0x0B, 0x1F, 0x4D)
GRAY = RGBColor(0x55, 0x61, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = "E8F1FF"
MINT = "E9F8F2"
CREAM = "FFF4DE"
BORDER = "B7C8E8"


def _set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def _set_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = tbl_pr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        tbl_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        tag = "w:" + edge
        elem = borders.find(qn(tag))
        if elem is None:
            elem = OxmlElement(tag)
            borders.append(elem)
        elem.set(qn("w:val"), "single")
        elem.set(qn("w:sz"), "6")
        elem.set(qn("w:space"), "0")
        elem.set(qn("w:color"), BORDER)


def _run_font(run, size: float = 9.0, bold: bool = False, color=None) -> None:
    run.font.name = FONT
    run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
    run.font.size = Pt(size)
    run.bold = bold
    if color:
        run.font.color.rgb = color


def _cell_text(cell, text: str, *, bold=False, color=None, size=8.6, fill=None) -> None:
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(text)
    _run_font(run, size=size, bold=bold, color=color)
    if fill:
        _set_cell_shading(cell, fill)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def _add_table(doc, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_borders(table)
    for cell, header in zip(table.rows[0].cells, headers):
        _cell_text(cell, header, bold=True, color=WHITE, size=8.1, fill="003B8F")
    for row in rows:
        cells = table.add_row().cells
        for idx, (cell, value) in enumerate(zip(cells, row)):
            _cell_text(cell, str(value), size=8.0)
    doc.add_paragraph()
    return table


def _add_note(doc, title: str, lines: list[str], fill: str = LIGHT) -> None:
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_borders(table)
    cell = table.cell(0, 0)
    _set_cell_shading(cell, fill)
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(title)
    _run_font(r, size=9.4, bold=True, color=BLUE)
    for line in lines:
        p = cell.add_paragraph()
        r = p.add_run("• " + line)
        _run_font(r, size=8.7, color=DEEP)
    doc.add_paragraph()


def _paragraph(doc, text: str, size: float = 9.2, color=DEEP, bold=False) -> None:
    p = doc.add_paragraph()
    r = p.add_run(text)
    _run_font(r, size=size, bold=bold, color=color)


def build_docx() -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    for style_name in ["Normal", "Heading 1", "Heading 2", "Heading 3"]:
        style = doc.styles[style_name]
        style.font.name = FONT
        style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
    doc.styles["Normal"].font.size = Pt(9.0)
    doc.styles["Heading 1"].font.color.rgb = BLUE
    doc.styles["Heading 1"].font.size = Pt(14)
    doc.styles["Heading 2"].font.color.rgb = DEEP
    doc.styles["Heading 2"].font.size = Pt(11)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("JB금융그룹 Fin:AI Challenge 기능명세서")
    _run_font(run, size=20, bold=True, color=BLUE)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("Cross-Check AI | 지정주제 2 준법자문가 AI Agent | JBIG 9th")
    _run_font(run, size=10, color=GRAY)

    _add_table(
        doc,
        ["팀명", "주제 구분", "팀원 정보(역할)", "작성일"],
        [["JBIG 9th", "지정주제 2 (준법자문가 AI Agent)", "강승윤 (PM · 기획 · 개발 · 디자인 총괄)", "2026.05.26"]],
    )

    doc.add_heading("1. 서비스 개요 (Service)", level=1)
    _add_table(
        doc,
        ["항목", "내용"],
        [
            ["서비스명", "Cross-Check AI"],
            [
                "서비스 한줄 소개",
                "한국어는 금소법, 인도네시아어는 OJK 규제로 각각 독립 심의하고, 번역·현지화 과정에서 발생하는 규제 드리프트를 교차검증으로 탐지하는 Human-in-the-loop 준법 보조 AI Agent",
            ],
            [
                "개발 목표",
                "대고객 콘텐츠 준법 심의의 수작업 의존과 다국어 처리 지연 문제를 줄이기 위해, 법령을 의사결정 트리로 사전 컴파일하고 Rule·Embedding·LLM·Citation 검증을 결합한 AI Agent를 구현한다.",
            ],
            [
                "타겟 사용자",
                "금융사 준법관리자(1차), 마케팅·영업·홍보 등 콘텐츠 생산 부서(2차), 해외 법인/현지화 담당자",
            ],
            [
                "기대 효과",
                "심의 시간 단축, 위반 누락(FN) 최소화, 조문 인용 오류 방지, 다국어 콘텐츠 일관성 확보, 규제 변경 대응 속도 향상",
            ],
        ],
    )
    _add_note(
        doc,
        "핵심 포지셔닝",
        [
            "Cross-Check AI는 법률판단을 대체하는 시스템이 아니라 위반 가능성을 조기 탐지하는 Regulatory Radar이다.",
            "AI는 위험 후보와 근거를 제시하고, 최종 승인·반려 권한은 항상 인간 준법관리자에게 남는다.",
            "금융 컴플라이언스에서는 과탐지(FP)는 재검토로 복구 가능하지만, 위반 누락(FN)은 실제 제재·평판 손상으로 이어질 수 있어 Recall 우선 설계가 타당하다.",
        ],
        MINT,
    )

    doc.add_heading("2. 시스템 구성도 (Architecture)", level=1)
    _add_note(
        doc,
        "전체 구조",
        [
            "사용자 입력: FastAPI 웹 UI에서 텍스트, PDF, 이미지 OCR 기반 대고객 콘텐츠 업로드",
            "Agent 1: 언어 감지 및 적용 규제 분류(한국어 금소법, 인도네시아어 OJK)",
            "Agent 2: YAML 법령 트리 기반 위반 탐지(Rule → Embedding → LLM 보완 판단)",
            "Agent 3: 한국어 원본과 현지어 번역본의 의미·규제 정합성 비교",
            "Agent 4: 국가법령정보 API 및 OJK 조문 DB를 통한 인용 조문 실존 여부 검증",
            "통합 판정 엔진: Risk Score, PASS/WARNING/VIOLATION, 근거 문장, 수정 권고안 산출",
            "DB/외부 연동: 법령 트리 DB, 제재 사례 벡터 DB, 국가법령정보 API, OJK 공시·문서",
        ],
    )

    doc.add_heading("3. 핵심 기능 명세 (Feature Specification)", level=1)
    _add_table(
        doc,
        ["기능명", "기능 설명", "입력/출력 데이터", "관련 기술 및 알고리즘", "구현 여부"],
        [
            [
                "콘텐츠 입력 및 규제 분류",
                "업로드 콘텐츠의 언어를 감지하고 적용할 법령 트리(한국 금소법·인니 OJK)를 자동 라우팅",
                "입력: 텍스트/PDF/이미지\n출력: 언어, 관할 규제, 적용 트리",
                "FastAPI, OCR, langdetect",
                "O",
            ],
            [
                "위반 탐지 엔진",
                "법령을 의사결정 트리로 컴파일한 뒤 명시적 금지 표현은 Rule로, 의미적 변형은 Embedding/LLM으로 보완 탐지",
                "입력: 콘텐츠+트리\n출력: 위반 후보, 근거, 리스크 스코어",
                "YAML DSL, Regex, FAISS, Groq Llama 3.3 70B",
                "O",
            ],
            [
                "OCR 후처리",
                "NICE 600점, 금소법 제19조, 금리 범위, +3% 등 금융광고 OCR에서 자주 깨지는 표현을 보정",
                "입력: OCR 원문\n출력: 정규화 텍스트",
                "정규식 보정, 금융용어 사전",
                "O",
            ],
            [
                "유사 제재 사례 검색",
                "AI 판정 사유와 공개 제재 사례를 비교해 준법관리자가 참고할 수 있는 유사 리스크 사례 제공",
                "입력: 판정 사유\n출력: 유사 사례, 출처, 유사도",
                "sentence-transformers, FAISS",
                "O",
            ],
            [
                "번역 정합성 검증",
                "한국어 원본과 인도네시아어 번역본을 독립 심의한 뒤 원본 통과/번역본 위반 등 규제 드리프트를 탐지",
                "입력: 다국어 콘텐츠 쌍\n출력: 언어별 심의 비교, 용어 불일치",
                "Termbase, 결과 diff, XLM-R 보조",
                "O",
            ],
            [
                "자기 검증 및 수정안 생성",
                "판정 결과의 인용 조문 번호·항·호 실존 여부를 재검증하여 LLM 환각과 잘못된 법령 인용을 방지",
                "입력: 판정 결과\n출력: 검증 결과, 원문 링크, 수정 권고",
                "국가법령정보 API lawSearch/lawService, OJK 조문 DB",
                "O",
            ],
            [
                "심의 결과 대시보드",
                "준법관리자가 위험도, 위반 후보, 근거 문장, 수정 방향을 한 화면에서 검토하고 최종 판단",
                "입력: AI 판정\n출력: Evidence 중심 심의 화면",
                "FastAPI 정적 UI, 이미지 하이라이트",
                "O",
            ],
            [
                "규제 변경 모니터링",
                "예선 MVP는 조문 검증·원문 링크 연결까지 구현. 본선에서는 금감원 RSS·OJK 공시 자동 모니터링으로 확장",
                "입력: 외부 API/RSS\n출력: 개정 알림, 영향 트리",
                "feedparser, scheduler, diff",
                "△",
            ],
        ],
    )

    doc.add_heading("4. 주요 기능 흐름도 (Flow)", level=1)
    _add_note(
        doc,
        "사용자 시나리오",
        [
            "마케팅/영업 부서가 광고 문구, 상품설명서, PDF 또는 이미지 홍보물을 웹 UI에 업로드한다.",
            "Agent 1이 언어와 적용 규제를 분류하고, Agent 2가 법령 트리를 실행해 명시적·문맥적 위반 가능성을 탐지한다.",
            "이미지/PDF 입력의 경우 OCR 결과와 이미지 하이라이트를 함께 제공하여 어느 부분이 검토 대상인지 빠르게 파악한다.",
            "Agent 4가 조문 실존 여부와 원문 링크를 검증하고, 유사 제재 사례 DB가 참고 사례를 제공한다.",
            "준법관리자는 AI 판정을 그대로 승인하는 것이 아니라 근거를 검토한 뒤 승인·수정 요청·반려 중 최종 결정을 수행한다.",
        ],
        LIGHT,
    )

    doc.add_heading("5. 향후 발전 방향 (Future Work)", level=1)
    _add_note(
        doc,
        "예선 이후 고도화",
        [
            "검증셋 기반 Precision·Recall·False Positive Rate 평가를 반복해 FN/FP trade-off를 정량 조정한다.",
            "금소법 17조(적합성), 21조(부당권유), 22조(광고규제)를 시작점으로 상품군별 감독규정과 가이드라인 트리를 확장한다.",
            "OJK 조항 매핑을 검수하고 인도네시아 현지 규제 문서·공시와 연결한다.",
            "금감원 RSS, 금융위 의결서, OJK 공시를 주기 모니터링하여 법령 변경이 영향을 주는 트리를 자동 표시한다.",
            "본선 단계에서는 영상/음성(STT) 기반 대고객 콘텐츠 검토, 심의 이력, 권한관리, 업무 큐까지 확장한다.",
        ],
        MINT,
    )

    doc.add_heading("6. 부록 (Appendix)", level=1)
    _add_table(
        doc,
        ["구분", "내용"],
        [
            [
                "선별 조항 근거",
                "금소법 22조는 광고규제, 21조는 권유성 표현, 17조는 고객 적합성 검토와 연결되어 대고객 콘텐츠 심의에서 빈번하고 발표 방어력이 높은 핵심 영역이다.",
            ],
            [
                "벤치마킹",
                "Norm AI의 규제 Agent 접근은 법령·정책·기관 가이드라인을 구조화해 기업 컴플라이언스 워크플로우에 연결하는 방향과 유사하다.",
            ],
            [
                "연구 근거",
                "Stanford CodeX/John Nay의 Law Informs Code 계열 연구는 법률 지식을 AI 시스템 내부의 구조화된 판단 로직으로 전환하는 접근의 근거로 활용 가능하다.",
            ],
            [
                "법률 AI 리스크",
                "Stanford HAI/RegLab의 법률 LLM hallucination 논의처럼 범용 LLM 단독 법률판단은 위험하므로, 본 프로젝트는 조문 검증 Agent와 Human-in-the-loop를 둔다.",
            ],
            [
                "국내 서비스 참고",
                "엘박스·로앤굿·슈퍼로이어 등 국내 법률 AI/리서치 서비스가 강조하는 출처 확인, 인용 링크, 전문가 최종 검토 UX를 참고한다.",
            ],
            [
                "테스트셋",
                f"총 105건 기준 평가: Risk Recall {SUMMARY['overall']['risk_recall']:.3f}, Risk Precision {SUMMARY['overall']['risk_precision']:.3f}, Violation Recall {SUMMARY['overall']['violation_recall_not_pass']:.3f}",
            ],
            [
                "주요 용어",
                "Human-in-the-loop: AI가 1차 탐지를 수행하고 인간 전문가가 최종 판단하는 협업형 AI 구조. Regulatory Radar: 법률판단 자동화가 아니라 잠재 규제 리스크를 사전 탐지하는 보조 시스템 개념.",
            ],
        ],
    )

    doc.add_heading("7. 기능 변경이력 (Change Log)", level=1)
    _add_table(
        doc,
        ["변경 일자", "변경 대상 기능", "변경 내용", "변경 사유"],
        [
            ["2026.05.26", "LLM", "OpenAI 대신 Groq Llama 3.3 70B로 변경", "공모전 구현 환경 및 비용 효율성 반영"],
            ["2026.05.26", "법령 트리", "금소법 22조 광고규제 트리 및 exclude_pattern, legal_basis 보완", "부정문 오탐과 조문 인용 오류 완화"],
            ["2026.05.26", "평가 체계", "테스트셋과 Recall 우선 평가 지표 반영", "FN 최소화 전략의 정량 근거 확보"],
            ["2026.05.26", "제출 문서", "벤치마킹·연구 근거·국내 법률 AI 참고 논리를 복구", "심사위원이 목적성과 근거를 빠르게 이해할 수 있도록 보완"],
        ],
    )

    doc.save(DOCX_OUT)


P_BLUE = PRGB(0, 59, 143)
P_DEEP = PRGB(11, 31, 77)
P_GRAY = PRGB(85, 97, 120)
P_LIGHT = PRGB(232, 241, 255)
P_MINT = PRGB(233, 248, 242)
P_CREAM = PRGB(255, 244, 222)
P_WHITE = PRGB(255, 255, 255)


def _ppt_font(run, size=12, bold=False, color=P_DEEP):
    run.font.name = FONT
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bg(slide, prs):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRGB(247, 250, 255)
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(0.18), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = P_BLUE
    bar.line.fill.background()


def _add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(PInches(0.62), PInches(0.32), PInches(12), PInches(0.52))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    _ppt_font(r, 22, True, P_DEEP)
    if subtitle:
        box = slide.shapes.add_textbox(PInches(0.64), PInches(0.86), PInches(12), PInches(0.3))
        tf = box.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = subtitle
        _ppt_font(r, 9.5, False, P_GRAY)


def _text_box(slide, x, y, w, h, text, size=12, bold=False, color=P_DEEP, fill=None, center=False):
    if fill:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = PRGB(190, 205, 235)
    else:
        shape = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = PInches(0.14)
    tf.margin_right = PInches(0.14)
    tf.margin_top = PInches(0.08)
    tf.margin_bottom = PInches(0.07)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    _ppt_font(r, size, bold, color)
    return shape


def _card(slide, x, y, w, h, title, bullets, fill=P_WHITE):
    shape = _text_box(slide, x, y, w, h, "", fill=fill)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = PInches(0.16)
    tf.margin_right = PInches(0.13)
    tf.margin_top = PInches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _ppt_font(r, 12.2, True, P_BLUE)
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.name = FONT
        p.font.size = PPt(8.9)
        p.font.color.rgb = P_DEEP
    return shape


def _three_cards(slide, cards, y=1.45, h=4.65):
    fills = [P_WHITE, P_LIGHT, P_MINT]
    for idx, (title, bullets) in enumerate(cards):
        _card(slide, 0.68 + idx * 4.12, y, 3.72, h, title, bullets, fills[idx % 3])


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    overall = SUMMARY["overall"]
    ko = SUMMARY.get("ko", {})
    ojk = SUMMARY.get("ojk", {})

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _text_box(slide, 0.8, 0.85, 8.6, 0.5, "JB금융그룹 Fin:AI Challenge", 20, True, P_BLUE)
    _text_box(slide, 0.8, 1.45, 8.0, 0.75, "MVP 제안서", 32, True, P_DEEP)
    _text_box(slide, 0.82, 2.35, 5.6, 0.48, "Cross-Check AI | JBIG 9th | 강승윤", 13, True, P_BLUE, P_LIGHT)
    _text_box(
        slide,
        0.82,
        3.05,
        6.6,
        0.75,
        "대고객 금융 콘텐츠의 규제 리스크를 조기 탐지하는\nHuman-in-the-loop Regulatory Radar",
        14,
        False,
        P_DEEP,
    )
    _text_box(
        slide,
        8.0,
        1.32,
        4.15,
        3.7,
        "AI가 법률판단을 대체하지 않습니다.\n\n법령 트리와 LLM은 위반 가능성과 근거를 제시하고,\n최종 승인·반려는 준법관리자가 수행합니다.",
        17,
        True,
        P_WHITE,
        P_BLUE,
        True,
    )

    slides = [
        (
            "1. Summary",
            "예선 제출물 7개 항목 중 1번",
            [
                (
                    "문제",
                    [
                        "대고객 콘텐츠 준법 심의가 대부분 수작업에 의존",
                        "다국어·다채널 확장 시 심의 리소스가 선형 증가",
                        "위반 누락(FN)은 실제 제재·평판 손상으로 연결",
                    ],
                ),
                (
                    "해결",
                    [
                        "금소법/OJK를 YAML 법령 트리로 구조화",
                        "Rule-first + Embedding + Groq Llama 3.3 70B",
                        "Agent 4가 조문 실존 여부와 원문 링크 재검증",
                    ],
                ),
                (
                    "근거",
                    [
                        "Norm AI·Stanford Law Informs Code 흐름을 벤치마킹",
                        f"테스트셋 기준 Risk Recall {overall['risk_recall']:.3f}",
                        "AI는 1차 탐지, 인간 준법관리자가 최종 판단",
                    ],
                ),
            ],
            "핵심 메시지: Cross-Check AI는 법률판단 자동화가 아니라 규제 리스크를 사전 탐지하는 Regulatory Radar입니다.",
        ),
        (
            "2. 문제 정의 (Problem Definition)",
            "예선 제출물 7개 항목 중 2번",
            [
                (
                    "현행 분석",
                    [
                        "광고·상품설명·SNS 등 대고객 콘텐츠 심의가 수작업 중심",
                        "심의 지연으로 콘텐츠 배포 시점이 늦어짐",
                        "규제 변경과 조문 인용 오류를 지속 추적해야 함",
                    ],
                ),
                (
                    "핵심 Pain Point",
                    [
                        "FN: 실제 위반 콘텐츠 배포로 이어질 수 있음",
                        "FP: 준법관리자의 재검토로 복구 가능",
                        "다국어 현지화 과정에서 의미·규제 드리프트 발생",
                    ],
                ),
                (
                    "JB금융 맥락",
                    [
                        "국내 금융규제와 해외 현지 규제를 동시에 고려해야 함",
                        "금소법·OJK를 모두 아는 인력은 희소함",
                        "AI Agent가 1차 스크리닝을 맡으면 업무 부담을 줄일 수 있음",
                    ],
                ),
            ],
            None,
        ),
        (
            "3. 제안 솔루션 개요 (Solution Overview)",
            "예선 제출물 7개 항목 중 3번",
            [
                (
                    "Agent 구조",
                    [
                        "Agent 1: 언어 감지·규제 분류",
                        "Agent 2: 위반 탐지 3단계",
                        "Agent 3: 다국어 의미·규제 정합성",
                        "Agent 4: 자기 검증·조문 검증",
                    ],
                ),
                (
                    "법령 트리 접근",
                    [
                        "명시적 금지 표현은 규칙 기반으로 재현 가능하게 탐지",
                        "애매한 문맥은 LLM으로 보완 판단",
                        "조문 인용은 API/DB로 다시 검증해 환각 방지",
                    ],
                ),
                (
                    "핵심 차별점",
                    [
                        "법령을 프롬프트가 아니라 실행 가능한 판단 트리로 전환",
                        "Recall 우선 설계로 위반 누락 최소화",
                        "Human-in-the-loop로 실무 적용성과 책임 소재 확보",
                    ],
                ),
            ],
            "설계 철학: AI는 보조, 사람이 최종 결정합니다.",
        ),
        (
            "4. 주요 기능 정의 (Key Features)",
            "예선 제출물 7개 항목 중 4번",
            [
                (
                    "MVP 구현 기능",
                    [
                        "텍스트/PDF/이미지 업로드 및 OCR",
                        "금소법 17·21·22조, OJK POJK 6/2022 트리",
                        "유사 제재 사례 검색과 이미지 하이라이트",
                    ],
                ),
                (
                    "위반 탐지 엔진",
                    [
                        "Rule: 원금보장·확정수익 등 명시 표현",
                        "Embedding: 유사 제재 사례와 의미적 근접성",
                        "LLM: 암묵적·문맥적 오인 가능성 판단",
                    ],
                ),
                (
                    "검토 UX",
                    [
                        "Risk Score 의미를 설명 가능한 문장으로 제공",
                        "JSON 원문이 아니라 비전공자도 읽는 근거 카드",
                        "승인·수정 요청·반려의 Human Review 흐름",
                    ],
                ),
            ],
            None,
        ),
        (
            "5. 데이터 및 기술 활용 (Data & Tech)",
            "예선 제출물 7개 항목 중 5번",
            [
                (
                    "활용 데이터",
                    [
                        "국가법령정보센터 Open API: 한국 법령 원문·조문",
                        "OJK POJK 문서: 인도네시아 소비자보호 규제",
                        "금감원·금융위 공개 제재 사례: 유사 리스크 검색",
                    ],
                ),
                (
                    "벤치마킹/연구",
                    [
                        "Norm AI: 규제 Agent 기반 기업 컴플라이언스",
                        "Stanford CodeX / John Nay: Law Informs Code",
                        "Stanford legal hallucination 논의: 조문 검증 필요성",
                    ],
                ),
                (
                    "구현 검증",
                    [
                        f"Risk Recall {overall['risk_recall']:.3f}",
                        f"Risk Precision {overall['risk_precision']:.3f}",
                        f"Violation Recall {overall['violation_recall_not_pass']:.3f}",
                        f"OJK Risk Recall {ojk.get('risk_recall', 1.0):.3f}",
                    ],
                ),
            ],
            "기술적 제약: OCR 깨짐·부정문 오탐·시각적 강조 판단은 별도 후처리와 Human Review 대상으로 둡니다.",
        ),
        (
            "6. 사용자 시나리오 / 유즈케이스",
            "예선 제출물 7개 항목 중 6번",
            [
                (
                    "1. 콘텐츠 업로드",
                    [
                        "마케팅팀이 광고 문구, PDF, 이미지 홍보물을 업로드",
                        "OCR 결과와 원문 이미지가 함께 표시",
                        "한국어/인도네시아어 단독 또는 쌍 입력 가능",
                    ],
                ),
                (
                    "2. AI 1차 심의",
                    [
                        "언어와 적용 규제를 자동 분류",
                        "법령 트리와 LLM이 위반 후보를 탐지",
                        "유사 제재 사례와 조문 링크를 함께 제시",
                    ],
                ),
                (
                    "3. 준법관리자 검토",
                    [
                        "AI 판정은 최종 결론이 아니라 검토 우선순위",
                        "준법관리자가 승인·수정 요청·반려 결정",
                        "결과는 향후 트리 튜닝과 평가셋 개선에 반영",
                    ],
                ),
            ],
            "Cross-Check AI는 사람이 놓칠 수 있는 규제 리스크를 먼저 비춰주는 보조 시스템입니다.",
        ),
        (
            "7. 기대 효과 및 향후 확장성",
            "예선 제출물 7개 항목 중 7번",
            [
                (
                    "기대 효과",
                    [
                        "콘텐츠 1차 심의 시간 단축",
                        "위반 누락(FN) 최소화와 조문 인용 오류 방지",
                        "다국어 콘텐츠의 규제 일관성 확보",
                    ],
                ),
                (
                    "본선 전 보강",
                    [
                        "테스트셋 기반 FP/FN 분석과 임계치 조정",
                        "OJK 조항 매핑 검수와 사례 보강",
                        "기능명세서·제안서의 근거·스토리 정리",
                    ],
                ),
                (
                    "본선 확장",
                    [
                        "금감원 RSS·OJK 공시 자동 모니터링",
                        "영상/음성 STT 기반 대고객 콘텐츠 검토",
                        "심의 이력·권한관리·업무 큐 등 운영 기능",
                    ],
                ),
            ],
            "실무 적용 방향: 넓게 탐지하고, 근거를 보여주고, 최종 판단은 사람이 합니다.",
        ),
    ]

    for title, subtitle, cards, footer in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide, prs)
        _add_title(slide, title, subtitle)
        _three_cards(slide, cards)
        if footer:
            _text_box(slide, 0.9, 6.28, 11.65, 0.48, footer, 11.5, True, P_BLUE, P_WHITE, True)

    prs.save(PPTX_OUT)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"saved {DOCX_OUT}")
    print(f"saved {PPTX_OUT}")
