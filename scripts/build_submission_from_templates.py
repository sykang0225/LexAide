import json
from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PPt


ROOT = Path(__file__).resolve().parents[2]
PROJECT = ROOT / "cross_check_ai"
SPEC_TEMPLATE = ROOT / "template_spec.docx"
PPT_TEMPLATE = ROOT / "template_mvp.pptx"
SPEC_OUT = ROOT / "CrossCheckAI_기능명세서.docx"
PPT_OUT = ROOT / "CrossCheckAI_MVP제안서.pptx"
SUMMARY = json.loads(
    (PROJECT / "data/evaluation/evaluation_report.summary.json").read_text(encoding="utf-8")
)

BLUE = RGBColor(0x00, 0x3B, 0x8F)
DEEP = RGBColor(0x0B, 0x1F, 0x4D)
BORDER = "B7C8E8"


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = tbl_pr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        tbl_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        tag = "w:" + edge
        elem = borders.find(qn(tag))
        if elem is None:
            elem = OxmlElement(tag)
            borders.append(elem)
        elem.set(qn("w:val"), "single")
        elem.set(qn("w:sz"), "6")
        elem.set(qn("w:space"), "0")
        elem.set(qn("w:color"), BORDER)


def set_cell(cell, text: str, bold=False, color=None, size=9, fill=None) -> None:
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color
    if fill:
        set_cell_shading(cell, fill)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def set_multiline_cell(cell, title: str, lines: list[str], fill=None) -> None:
    cell.text = ""
    if fill:
        set_cell_shading(cell, fill)
    paragraph = cell.paragraphs[0]
    title_run = paragraph.add_run(title)
    title_run.bold = True
    title_run.font.color.rgb = BLUE
    title_run.font.name = "맑은 고딕"
    title_run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    title_run.font.size = Pt(9.2)
    for line in lines:
        paragraph = cell.add_paragraph()
        run = paragraph.add_run("· " + line)
        run.font.name = "맑은 고딕"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
        run.font.size = Pt(8.8)


def build_docx() -> None:
    doc = Document(SPEC_TEMPLATE)

    for style_name in ["Normal", "Heading 1", "Heading 2", "Heading 3"]:
        if style_name in doc.styles:
            style = doc.styles[style_name]
            style.font.name = "맑은 고딕"
            style._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")

    # 0. Header metadata table
    meta = doc.tables[0]
    set_borders(meta)
    set_cell(meta.cell(1, 0), "팀명", True, BLUE, 9, "E8F1FF")
    set_cell(meta.cell(1, 1), "JBIG 9th", False, None, 9)
    set_cell(meta.cell(1, 2), "주제 구분", True, BLUE, 9, "E8F1FF")
    set_cell(meta.cell(1, 3), "지정주제 2 (준법자문가 AI Agent)", False, None, 9)
    set_cell(meta.cell(2, 0), "팀원 정보(역할)", True, BLUE, 9, "E8F1FF")
    set_cell(meta.cell(2, 1), "강승윤 (PM · 기획 · 개발 · 디자인 총괄)", False, None, 9)
    set_cell(meta.cell(2, 2), "작성일", True, BLUE, 9, "E8F1FF")
    set_cell(meta.cell(2, 3), "2026.05.26", False, None, 9)

    # 1. Service overview
    overview = doc.tables[1]
    set_borders(overview)
    overview_rows = {
        "서비스명": "Cross-Check AI",
        "서비스 한줄 소개": "대고객 금융 콘텐츠의 규제 리스크를 법령 트리와 LLM으로 1차 탐지하고, 준법관리자가 최종 승인하는 Human-in-the-loop Regulatory Radar",
        "개발 목표": "수작업 중심 준법 심의의 지연·품질 편차·휴먼에러를 줄이고, 한국 금소법과 인도네시아 OJK 규정을 함께 검토할 수 있는 AI Agent 서비스를 구현합니다.",
        "타겟 사용자": "금융사 준법관리자, 마케팅·영업·홍보 콘텐츠 제작 부서, 해외 법인 현지화 담당자",
        "기대 효과": "심의 리드타임 단축, 위반 누락(FN) 최소화, 조문 인용 오류 방지, 다국어 콘텐츠의 현지 규제 대응력 확보",
    }
    for row in overview.rows:
        key = row.cells[0].text.strip()
        if key in overview_rows:
            set_cell(row.cells[0], key, True, BLUE, 9, "E8F1FF")
            set_cell(row.cells[1], overview_rows[key], False, None, 8.8)

    # 2. Architecture
    arch = doc.tables[2]
    set_borders(arch)
    set_multiline_cell(
        arch.cell(0, 0),
        "① 사용자 입력 → ② AI Agent → ③ 백엔드/API → ④ DB/외부연동 → ⑤ 프론트엔드/UI",
        [
            "사용자 입력: 텍스트, PDF, 이미지 OCR 기반 대고객 콘텐츠 업로드",
            "Agent 1: 언어 감지 및 금소법/OJK 규제 트리 라우팅",
            "Agent 2: Rule → Embedding → LLM 하이브리드 위반 탐지",
            "Agent 3: 한국어 원본과 인도네시아어 현지본의 의미·규제 정합성 비교",
            "Agent 4: 국가법령정보 API 및 OJK 조문 DB 기반 인용 조문 실존 검증",
            "FastAPI 웹앱: Evidence 중심 결과 화면과 준법관리자 승인/수정/반려 흐름 제공",
        ],
        "F5F8FF",
    )

    # 3. Feature specification
    features = doc.tables[3]
    set_borders(features)
    while len(features.rows) < 8:
        features.add_row()
    headers = ["기능명", "기능 설명", "입력/출력 데이터", "관련 기술 및 알고리즘", "구현 여부"]
    for cell, header in zip(features.rows[0].cells, headers):
        set_cell(cell, header, True, RGBColor(255, 255, 255), 8.2, "003B8F")
    rows = [
        ["콘텐츠 입력 및 규제 분류", "언어를 감지하고 한국 금소법 또는 인도네시아 OJK 트리로 자동 라우팅", "입력: 텍스트/PDF/이미지\n출력: 언어·규제 범위", "langdetect, OCR, FastAPI", "O"],
        ["위반 탐지 엔진", "법령을 YAML 의사결정 트리로 컴파일하고 명시적 위반은 Rule로 우선 탐지", "입력: 콘텐츠+트리\n출력: PASS/WARNING/VIOLATION", "YAML DSL, Regex, Groq Llama 3.3 70B", "O"],
        ["근거 조문 검증", "탐지 결과의 인용 조문 실존 여부를 검증해 LLM 환각과 잘못된 인용 방지", "입력: citation\n출력: 검증 결과·원문 링크", "국가법령정보 API lawSearch/lawService, OJK 조문 DB", "O"],
        ["유사 제재 사례 검색", "위반 사유와 유사한 공개 제재 사례를 검색해 준법관리자 판단 보조", "입력: 판정 결과\n출력: 유사 사례", "FAISS, sentence-transformers, 제재 사례 DB", "O"],
        ["다국어 정합성 검증", "한국어 원본과 인니어 현지본을 독립 심의 후 결과 차이를 비교", "입력: 원본+현지본\n출력: 규제 드리프트", "Termbase, 결과 diff, 교차검증 Agent", "O"],
        ["심의 결과 UI", "법률 JSON 대신 Evidence·Risk Score·수정 권고·Human Review 버튼 제공", "입력: AI 판정\n출력: 준법 심의 화면", "FastAPI 정적 웹 UI", "O"],
        ["규제 모니터링", "예선은 조문 검증까지 구현, 본선에서 RSS/OJK 공시 자동 모니터링으로 확장", "입력: API/RSS\n출력: 개정 알림", "feedparser, law API, scheduler", "△"],
    ]
    for idx, values in enumerate(rows, start=1):
        for cell, value in zip(features.rows[idx].cells, values):
            set_cell(cell, value, value == "O", BLUE if value in {"O", "△"} else None, 7.6)

    # 4. Flow
    flow = doc.tables[4]
    set_borders(flow)
    set_multiline_cell(
        flow.cell(0, 0),
        "주요 기능 흐름",
        [
            "콘텐츠 업로드: 마케팅/영업 부서가 광고 문구, PDF, 이미지 홍보물을 업로드",
            "언어·규제 분류: Agent 1이 한국어/인도네시아어를 감지하고 금소법/OJK 트리를 선택",
            "위반 탐지: Agent 2가 명시적 표현은 Rule로, 애매한 문맥은 LLM으로 보완 검토",
            "근거 보강: 유사 제재 사례 검색과 Agent 4 조문 실존 검증으로 판단 근거를 정리",
            "사람 검토: 준법관리자가 AI 결과를 확인하고 승인·수정 요청·반려 중 최종 결정",
        ],
        "F5F8FF",
    )

    # 5. Future work
    future = doc.tables[5]
    set_borders(future)
    set_multiline_cell(
        future.cell(0, 0),
        "향후 발전 방향",
        [
            "평가셋 기반 FN/FP 리뷰를 통해 금소법·OJK 트리 정밀화",
            "제19조 설명의무, 제18조 적정성 원칙, 상품군별 감독규정으로 법령 범위 확장",
            "금감원 RSS와 OJK 공시 자동 모니터링을 통해 개정 법령 영향도 분석",
            "영상·음성 광고의 STT 분석 기능은 본선 단계에서 구현",
            "운영 피드백을 법령 트리와 제재 사례 DB에 반영하는 지속 개선 구조 도입",
        ],
        "F5F8FF",
    )

    # 6. Appendix
    appendix = doc.tables[6]
    set_borders(appendix)
    set_multiline_cell(
        appendix.cell(0, 0),
        "부록 및 참고자료",
        [
            "대회 상세주제 안내: 지정주제 2 준법자문가 AI Agent 서비스 개발",
            "법령 출처: 국가법령정보센터 Open API, POJK No. 6/POJK.07/2022",
            "평가셋: 총 105건(한국어 60건, 금소법 17조 15건, OJK 30건)",
            f"검증 결과: Risk Recall {SUMMARY['overall']['risk_recall']:.3f}, Risk Precision {SUMMARY['overall']['risk_precision']:.3f}, Violation Recall {SUMMARY['overall']['violation_recall_not_pass']:.3f}",
            "핵심 해석: Risk Score와 LLM Confidence는 최종 법률 판단이 아니라 준법관리자 검토 우선순위를 정하는 보조 지표",
        ],
        "F5F8FF",
    )

    # 7. Change log
    change = doc.tables[7]
    set_borders(change)
    for cell, header in zip(change.rows[0].cells, ["변경 일자", "변경 대상 기능", "변경 내용", "변경 사유"]):
        set_cell(cell, header, True, RGBColor(255, 255, 255), 8, "003B8F")
    change_rows = [
        ["2026.05.26", "법령 트리", "금소법 22조 및 OJK WARNING 룰 보강", "Risk Recall 개선 및 FN 감소"],
        ["2026.05.26", "평가 리포트", "105건 테스트셋 기준 평가 결과 반영", "Recall 우선 설계 근거 확보"],
        ["2026.05.26", "제출 문서", "공식 7개 항목 양식 유지 방식으로 재정리", "제출 형식 준수"],
    ]
    for idx, values in enumerate(change_rows, start=1):
        if idx >= len(change.rows):
            change.add_row()
        for cell, value in zip(change.rows[idx].cells, values):
            set_cell(cell, value, size=7.8)

    doc.save(SPEC_OUT)


P_BLUE = PRGB(0, 59, 143)
P_DEEP = PRGB(11, 31, 77)
P_GRAY = PRGB(82, 94, 120)
P_LIGHT = PRGB(232, 241, 255)
P_MINT = PRGB(232, 248, 242)
P_WHITE = PRGB(255, 255, 255)


def clear_instruction_shapes(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and "작성방법" in shape.text:
            shape.text_frame.clear()


def text_box(slide, x, y, w, h, text, size=12, bold=False, color=P_DEEP, fill=None, align="left"):
    if fill:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = PRGB(190, 206, 235)
    else:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.06)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def bullet_card(slide, x, y, w, h, title, bullets, fill=P_WHITE):
    text_box(slide, x, y, w, h, "", fill=fill)
    shape = slide.shapes[-1]
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.15)
    frame.margin_top = Inches(0.12)
    title_p = frame.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = title
    title_run.font.name = "맑은 고딕"
    title_run.font.size = PPt(12.5)
    title_run.font.bold = True
    title_run.font.color.rgb = P_BLUE
    for bullet in bullets:
        p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "맑은 고딕"
        p.font.size = PPt(9.7)
        p.font.color.rgb = P_DEEP


def add_metric(slide, x, label, value, note):
    text_box(slide, x, 4.55, 2.4, 1.05, "", fill=P_WHITE, align="center")
    shape = slide.shapes[-1]
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "맑은 고딕"
    r.font.size = PPt(21)
    r.font.bold = True
    r.font.color.rgb = P_BLUE
    p = frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "맑은 고딕"
    r.font.size = PPt(8.8)
    r.font.bold = True
    r.font.color.rgb = P_DEEP
    p = frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = note
    r.font.name = "맑은 고딕"
    r.font.size = PPt(7.2)
    r.font.color.rgb = P_GRAY


def build_pptx() -> None:
    prs = Presentation(PPT_TEMPLATE)
    overall = SUMMARY["overall"]
    ko = SUMMARY["ko"]
    ojk = SUMMARY["ojk"]

    # Slide 1: cover
    slide = prs.slides[0]
    text_box(slide, 0.8, 4.85, 5.3, 0.45, "Cross-Check AI · JBIG 9th · 강승윤", 14, True, P_BLUE, P_LIGHT)
    text_box(slide, 0.8, 5.45, 8.4, 0.42, "대고객 금융 콘텐츠 준법 심의를 위한 Human-in-the-loop AI Agent", 12, False, P_DEEP)

    # Slide 2: Summary
    slide = prs.slides[1]
    clear_instruction_shapes(slide)
    bullet_card(slide, 0.75, 1.45, 3.7, 3.55, "문제", [
        "대고객 콘텐츠 준법 심의가 대부분 수작업",
        "다국어 콘텐츠 확장 시 심의 리소스가 선형 증가",
        "규제 변경·조문 인용 오류·휴먼에러 가능성 존재",
    ], P_WHITE)
    bullet_card(slide, 4.8, 1.45, 3.7, 3.55, "해결", [
        "금소법·OJK 규정을 YAML 법령 트리로 구조화",
        "Rule-first + Groq Llama 3.3 70B + Agent 4 조문 검증",
        "AI는 1차 탐지, 최종 판단은 준법관리자",
    ], P_LIGHT)
    bullet_card(slide, 8.85, 1.45, 3.7, 3.55, "검증", [
        "105건 테스트셋 기준 Risk Recall 0.886",
        "Violation Recall 0.943, 한국어 위반 Recall 1.000",
        "FastAPI 웹앱, OCR, 유사 제재 검색, 조문 원문 링크 구현",
    ], P_MINT)
    text_box(slide, 1.05, 5.55, 11.2, 0.62, "핵심 메시지: Cross-Check AI는 법률판단을 대체하지 않고, 위반 가능성을 조기에 탐지하는 Regulatory Radar입니다.", 13, True, P_BLUE, P_WHITE, "center")

    # Slide 3: Problem Definition
    slide = prs.slides[2]
    clear_instruction_shapes(slide)
    bullet_card(slide, 0.8, 1.45, 3.65, 4.3, "현행 분석", [
        "준법관리자가 광고·상품설명·SNS 등 대고객 콘텐츠를 수작업 검토",
        "외국어 콘텐츠도 동일 인력이 동일 방식으로 심의하는 경우가 많음",
        "심의 지연으로 콘텐츠 배포 시점이 늦어질 수 있음",
    ])
    bullet_card(slide, 4.8, 1.45, 3.65, 4.3, "핵심 Pain Point", [
        "위반 누락(FN)은 실제 규제 위반 콘텐츠 배포로 이어질 수 있음",
        "규제 변경과 조문 인용 오류를 사람이 계속 추적해야 함",
        "다국어·다채널 확장 시 품질 편차와 휴먼에러 가능성 증가",
    ], P_LIGHT)
    bullet_card(slide, 8.8, 1.45, 3.65, 4.3, "해결 대상", [
        "1차 사용자: 금융사 준법관리자",
        "2차 사용자: 마케팅·영업·홍보 콘텐츠 제작 부서",
        "JB금융 인도네시아 등 해외 사업 확장에 따른 현지 규제 대응",
    ], P_MINT)

    # Slide 4: Solution Overview
    slide = prs.slides[3]
    clear_instruction_shapes(slide)
    text_box(slide, 0.85, 1.25, 11.65, 0.55, "Cross-Check AI = 법령 트리 기반 준법 심의 Agent + Human-in-the-loop 검토 콘솔", 15, True, P_BLUE, P_LIGHT, "center")
    for idx, (title, body) in enumerate([
        ("Agent 1", "언어 감지·규제 분류"),
        ("Agent 2", "Rule → Embedding → LLM 위반 탐지"),
        ("Agent 3", "다국어 의미·규제 정합성 비교"),
        ("Agent 4", "조문 실존 여부·원문 링크 검증"),
    ]):
        text_box(slide, 0.9 + idx * 3.05, 2.2, 2.45, 1.15, title + "\n" + body, 10.5, True, P_DEEP, P_WHITE, "center")
    bullet_card(slide, 0.95, 4.05, 3.55, 1.65, "설계 철학", ["AI는 법률판단 자동화가 아니라 준법관리자의 검토 효율화를 위한 보조 시스템"], P_WHITE)
    bullet_card(slide, 4.9, 4.05, 3.55, 1.65, "Recall 우선", ["과탐지는 재검토 가능하지만 위반 누락은 과징금·제재·평판 손상으로 연결 가능"], P_LIGHT)
    bullet_card(slide, 8.85, 4.05, 3.55, 1.65, "설명 가능성", ["위반 사유, 매칭 문구, 근거 조문, 유사 제재 사례를 함께 제시"], P_MINT)

    # Slide 5: Key Features
    slide = prs.slides[4]
    clear_instruction_shapes(slide)
    feature_cards = [
        ("1. 콘텐츠 입력·분류", ["텍스트/PDF/이미지 OCR", "한국어·인니어 자동 라우팅"]),
        ("2. 위반 탐지 엔진", ["금소법 17·21·22 + OJK 트리", "명시 위반은 Rule로 우선 탐지"]),
        ("3. 근거 조문 검증", ["국가법령정보 API lawSearch/lawService", "OJK 조문 DB로 citation 검증"]),
        ("4. 유사 제재 검색", ["공개 제재 사례 벡터 DB", "판단 근거 보강"]),
        ("5. 다국어 정합성", ["원본·현지본 독립 심의", "번역상 규제 드리프트 탐지"]),
        ("6. Human Review UI", ["Risk Score 설명", "승인·수정요청·반려 흐름"]),
    ]
    for idx, (title, bullets) in enumerate(feature_cards):
        x = 0.75 + (idx % 3) * 4.05
        y = 1.28 + (idx // 3) * 2.2
        bullet_card(slide, x, y, 3.55, 1.75, title, bullets, P_LIGHT if idx % 2 else P_WHITE)
    text_box(slide, 1.0, 5.95, 11.0, 0.55, "MVP 구현 완료: FastAPI 웹앱 · OCR · YAML 트리 엔진 · Agent 4 조문 검증 · 평가셋 리포트", 12, True, P_BLUE, P_MINT, "center")

    # Slide 6: Data & Tech
    slide = prs.slides[5]
    clear_instruction_shapes(slide)
    bullet_card(slide, 0.75, 1.35, 3.6, 4.25, "활용 데이터", [
        "국가법령정보센터 Open API",
        "POJK No. 6/POJK.07/2022 조문 DB",
        "공개 제재 사례 DB + FAISS 벡터",
        "105건 테스트셋(PASS/WARNING/VIOLATION)",
    ])
    bullet_card(slide, 4.85, 1.35, 3.6, 4.25, "기술 구조", [
        "FastAPI 정적 웹앱",
        "YAML DSL Decision Tree",
        "Groq Llama 3.3 70B",
        "Tesseract OCR + OCR 후처리",
        "sentence-transformers + FAISS",
    ], P_LIGHT)
    bullet_card(slide, 8.95, 1.35, 3.6, 4.25, "검증 결과", [
        f"Risk Recall {overall['risk_recall']:.3f}",
        f"Risk Precision {overall['risk_precision']:.3f}",
        f"Violation Recall {overall['violation_recall_not_pass']:.3f}",
        f"OJK Risk Recall {ojk['risk_recall']:.3f}",
    ], P_MINT)

    # Slide 7: User Scenario
    slide = prs.slides[6]
    clear_instruction_shapes(slide)
    steps = [
        ("1", "콘텐츠 업로드", "마케팅팀이 광고 문구·PDF·이미지 홍보물을 업로드"),
        ("2", "AI 1차 심의", "언어 감지 후 금소법/OJK 트리로 위반 가능성 탐지"),
        ("3", "근거 확인", "조문 원문 링크, 유사 제재 사례, OCR 품질을 함께 확인"),
        ("4", "사람 최종 결정", "준법관리자가 승인·수정 요청·반려 중 결정"),
    ]
    for idx, (num, title, body) in enumerate(steps):
        y = 1.35 + idx * 1.05
        text_box(slide, 0.95, y, 0.65, 0.62, num, 16, True, P_WHITE, P_BLUE, "center")
        text_box(slide, 1.85, y, 2.65, 0.62, title, 12, True, P_DEEP, P_LIGHT, "center")
        text_box(slide, 4.75, y, 7.4, 0.62, body, 11.2, False, P_DEEP, P_WHITE, "center")
    text_box(slide, 1.1, 5.9, 11.0, 0.55, "사용자 경험의 핵심은 법률 JSON을 보여주는 것이 아니라, 준법관리자가 바로 판단 가능한 Evidence 중심 화면을 제공하는 것입니다.", 12, True, P_BLUE, P_MINT, "center")

    # Slide 8: Expected Impact
    slide = prs.slides[7]
    clear_instruction_shapes(slide)
    bullet_card(slide, 0.75, 1.35, 3.6, 4.35, "예선 MVP 효과", [
        "심의 리스크를 사람이 보기 전에 1차 선별",
        "조문 인용 오류와 LLM 환각 방지",
        "테스트셋 기반 Recall 우선 설계 근거 확보",
    ])
    bullet_card(slide, 4.85, 1.35, 3.6, 4.35, "본선 확장", [
        "금감원 RSS·OJK 공시 자동 모니터링",
        "영상·음성 광고 STT 분석",
        "심의 이력, 권한, 업무 큐 기능",
    ], P_LIGHT)
    bullet_card(slide, 8.95, 1.35, 3.6, 4.35, "실무 적용성", [
        "준법관리자는 최종 판단에 집중",
        "운영 피드백으로 법령 트리 지속 개선",
        "국가·상품군별 규제 확장 가능",
    ], P_MINT)
    text_box(slide, 1.0, 5.95, 11.0, 0.55, "Cross-Check AI는 AI Agent 흐름에 맞춰 JB금융의 동남아·다국어 콘텐츠 준법 심의를 수동 운영보다 빠르고 일관되게 보조합니다.", 12, True, P_BLUE, P_WHITE, "center")

    prs.save(PPT_OUT)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"saved {SPEC_OUT}")
    print(f"saved {PPT_OUT}")
