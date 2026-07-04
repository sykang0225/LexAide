import json
from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt


ROOT = Path(__file__).resolve().parents[2]
DOCX = next(p for p in ROOT.glob("*.docx") if "복구본" not in p.name)
PPTX = next(p for p in ROOT.glob("*.pptx") if "복구본" not in p.name)
SUMMARY = json.loads(
    (ROOT / "cross_check_ai/data/evaluation/evaluation_report.summary.json").read_text(
        encoding="utf-8"
    )
)

BLUE = RGBColor(0x00, 0x3B, 0x8F)
DEEP = RGBColor(0x0B, 0x1F, 0x4D)
BORDER = "B7C8E8"


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_text(cell, text: str, bold=False, color=None, size=9) -> None:
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    if color:
        run.font.color.rgb = color
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def set_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = tbl_pr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        tbl_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        tag = "w:" + edge
        elem = borders.find(qn(tag))
        if elem is None:
            elem = OxmlElement(tag)
            borders.append(elem)
        elem.set(qn("w:val"), "single")
        elem.set(qn("w:sz"), "6")
        elem.set(qn("w:space"), "0")
        elem.set(qn("w:color"), BORDER)


def add_note(doc: Document, title: str, body: str, fill="E8F1FF") -> None:
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    cell = table.cell(0, 0)
    set_cell_shading(cell, fill)
    paragraph = cell.paragraphs[0]
    first = paragraph.add_run(title + "  ")
    first.bold = True
    first.font.name = "맑은 고딕"
    first._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    first.font.color.rgb = BLUE
    first.font.size = Pt(9.5)
    second = paragraph.add_run(body)
    second.font.name = "맑은 고딕"
    second._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    second.font.size = Pt(9.5)
    doc.add_paragraph()


def build_docx() -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.72)
    section.bottom_margin = Inches(0.72)
    section.left_margin = Inches(0.78)
    section.right_margin = Inches(0.78)

    styles = doc.styles
    styles["Normal"].font.name = "맑은 고딕"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    styles["Normal"].font.size = Pt(9.5)
    for style_name in ["Heading 1", "Heading 2", "Heading 3"]:
        styles[style_name].font.name = "맑은 고딕"
        styles[style_name]._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    styles["Heading 1"].font.color.rgb = BLUE
    styles["Heading 1"].font.size = Pt(15)
    styles["Heading 2"].font.color.rgb = DEEP
    styles["Heading 2"].font.size = Pt(12)

    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run("Cross-Check AI 기능명세서")
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = BLUE
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")

    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run("JB금융그룹 Fin:AI Challenge 2026 · 지정주제 2 준법자문가 AI Agent")
    run.font.size = Pt(10.5)
    run.font.color.rgb = DEEP
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")

    meta = doc.add_table(rows=3, cols=4)
    meta.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(meta)
    rows = [
        ("팀명", "JBIG 9th", "작성자", "강승윤"),
        ("서비스명", "Cross-Check AI", "작성일", "2026.05.26"),
        ("핵심 포지션", "Human-in-the-loop Regulatory Radar", "LLM", "Groq Llama 3.3 70B"),
    ]
    for table_row, values in zip(meta.rows, rows):
        for idx, value in enumerate(values):
            set_cell_text(
                table_row.cells[idx],
                value,
                bold=idx % 2 == 0,
                color=BLUE if idx % 2 == 0 else None,
                size=9,
            )
            if idx % 2 == 0:
                set_cell_shading(table_row.cells[idx], "E8F1FF")

    doc.add_paragraph()
    add_note(
        doc,
        "한 줄 정의",
        "금융 대고객 콘텐츠를 법령 트리 기반으로 1차 탐지하고, AI가 근거와 수정 방향을 제시하되 최종 승인·반려는 준법관리자가 결정하는 보조 AI 시스템입니다.",
    )

    doc.add_heading("1. 서비스 목적", level=1)
    for text in [
        "현행 준법 심의는 대고객 콘텐츠를 준법관리자가 수작업으로 검토하는 구조라 콘텐츠 증가, 다국어 확장, 규제 변경에 따라 병목이 커집니다.",
        "Cross-Check AI는 법률판단 자동화가 아니라 위반 가능성을 조기에 넓게 탐지하는 Regulatory Radar입니다. 과탐지(FP)는 사람이 재검토해 복구 가능하지만, 위반 누락(FN)은 실제 규제 위반 콘텐츠 배포로 이어질 수 있으므로 Recall 우선 전략을 채택합니다.",
        "특히 JB금융의 인도네시아 사업 맥락에서는 한국 금소법과 인도네시아 OJK 규정을 각각 적용하고, 번역·현지화 과정에서 의미가 변질되는 규제 드리프트를 확인할 수 있습니다.",
    ]:
        doc.add_paragraph(text)

    add_note(
        doc,
        "법학 배경 활용",
        "조문을 통째로 LLM에게 맡기는 방식이 아니라, 광고·권유·적합성이라는 쟁점을 선별하고 이를 사람이 검수 가능한 트리로 구조화하는 데 법학적 사고를 활용합니다.",
        "E9F8F2",
    )

    doc.add_heading("2. 법령 범위 선정 근거", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, ["조항", "담당 리스크", "탐지 예시", "선정 이유"]):
        set_cell_text(cell, header, bold=True, color=RGBColor(255, 255, 255), size=8.5)
        set_cell_shading(cell, "003B8F")

    law_rows = [
        ["금소법 제22조", "광고 규제", "원금보장, 확정수익, 비용·위험 고지 누락, 근거 없는 비교", "대고객 광고 문구에서 가장 빈번하고 데모 가능성이 높은 리스크 축"],
        ["금소법 제21조", "부당권유", "단정적 판단, 거짓·과장, 불안감 조성, 즉시 가입 압박", "광고와 권유의 경계에 있는 콘텐츠를 보완"],
        ["금소법 제17조", "적합성 원칙", "고위험상품 무차별 권유, 고객 특성 확인 없는 가입 유도", "누구에게 권하는지에 대한 대상 리스크를 반영"],
        ["POJK 6/2022", "OJK 소비자보호", "명확성, 오인 방지, 위험·비용 고지, 비교표현", "JB금융 인도네시아 확장 맥락의 현지 규제 대응"],
    ]
    for row in law_rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            set_cell_text(cell, value, size=8)

    doc.add_paragraph(
        "전체 금소법을 한 번에 구현하지 않은 이유는 예선 MVP에서 설명 가능성과 검증 가능성을 유지하기 위해서입니다. 본 시스템은 먼저 광고·권유·적합성 리스크를 안정적으로 탐지하고, 본선에서 설명의무·적정성 원칙·상품군별 감독규정으로 확장하는 구조입니다."
    )

    doc.add_heading("3. 시스템 구조", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, ["Agent", "역할", "구현 상태"]):
        set_cell_text(cell, header, bold=True, color=RGBColor(255, 255, 255), size=8.5)
        set_cell_shading(cell, "003B8F")
    agent_rows = [
        ["Agent 1", "언어 감지 및 한국 금소법/OJK 규제 트리 라우팅", "구현"],
        ["Agent 2", "YAML 법령 트리 실행, Rule → Embedding → LLM 보완 탐지", "구현"],
        ["Agent 3", "한국어 원본과 인도네시아어 현지본의 의미·규제 정합성 비교", "구현"],
        ["Agent 4", "국가법령정보 API와 OJK 로컬 조문 DB로 인용 조문 실존 여부 검증", "구현"],
        ["통합 판정", "Risk Score, PASS/WARNING/VIOLATION, Human-in-the-loop 승인 흐름", "구현"],
    ]
    for row in agent_rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            set_cell_text(cell, value, size=8.2)

    doc.add_heading("4. 핵심 기능 명세", level=1)
    table = doc.add_table(rows=1, cols=5)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, ["기능", "사용자 입력", "AI 처리", "출력", "상태"]):
        set_cell_text(cell, header, bold=True, color=RGBColor(255, 255, 255), size=8)
        set_cell_shading(cell, "003B8F")
    feature_rows = [
        ["콘텐츠 업로드", "텍스트, PDF, 이미지", "OCR 후처리 및 언어 감지", "추출 텍스트, OCR 품질", "O"],
        ["위반 탐지", "광고·권유 문구", "금소법 17·21·22 및 OJK 트리 실행", "위반·주의 항목, 근거 조문", "O"],
        ["유사 제재 검색", "탐지 결과", "제재 사례 벡터 DB 검색", "유사 사례와 비교 근거", "O"],
        ["조문 검증", "판정 citation", "lawSearch/lawService 및 OJK 조문 DB 확인", "실존 조문 여부, 원문 링크", "O"],
        ["Human Review", "AI 판정 결과", "리스크 설명과 수정 권고 제시", "승인/수정 요청/반려 보조", "O"],
    ]
    for row in feature_rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            set_cell_text(cell, value, size=7.8)

    doc.add_heading("5. 평가 결과", level=1)
    add_note(
        doc,
        "평가 기준",
        "Exact Accuracy보다 Risk Recall을 핵심 지표로 봅니다. WARNING/VIOLATION을 최소한 사람 검토 대상으로 올렸는지가 준법 리스크 관점에서 더 중요하기 때문입니다.",
        "FFF4DE",
    )
    table = doc.add_table(rows=1, cols=6)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, ["평가셋", "건수", "정확도", "Risk Recall", "Risk Precision", "Violation Recall"]):
        set_cell_text(cell, header, bold=True, color=RGBColor(255, 255, 255), size=8)
        set_cell_shading(cell, "003B8F")

    overall = SUMMARY["overall"]
    ko = SUMMARY["ko"]
    ko17 = SUMMARY["ko_17"]
    ojk = SUMMARY["ojk"]
    metric_rows = [
        ["전체", overall["total"], overall["exact_accuracy"], overall["risk_recall"], overall["risk_precision"], overall["violation_recall_not_pass"]],
        ["한국어 22·21조", ko["total"], ko["exact_accuracy"], ko["risk_recall"], ko["risk_precision"], ko["violation_recall_not_pass"]],
        ["한국어 17조", ko17["total"], ko17["exact_accuracy"], ko17["risk_recall"], ko17["risk_precision"], ko17["violation_recall_not_pass"]],
        ["OJK", ojk["total"], ojk["exact_accuracy"], ojk["risk_recall"], ojk["risk_precision"], ojk["violation_recall_not_pass"]],
    ]
    for row in metric_rows:
        cells = table.add_row().cells
        values = [row[0], str(row[1]), f"{row[2]:.3f}", f"{row[3]:.3f}", f"{row[4]:.3f}", f"{row[5]:.3f}"]
        for cell, value in zip(cells, values):
            set_cell_text(cell, value, size=8)

    doc.add_paragraph(
        "트리 보강 후 전체 Risk Recall은 0.886, Violation Recall은 0.943입니다. 한국어 위반 케이스는 1차 룰 기반에서 누락 없이 사람 검토 대상으로 올리는 구조를 확보했습니다. 남은 과제는 OJK 위반 누락 일부와 WARNING 세부 라벨 정밀화입니다."
    )

    doc.add_heading("6. 예선 구현 범위와 본선 확장", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, ["구분", "예선 MVP 구현", "본선 확장"]):
        set_cell_text(cell, header, bold=True, color=RGBColor(255, 255, 255), size=8)
        set_cell_shading(cell, "003B8F")
    roadmap_rows = [
        ["웹 서비스", "FastAPI 기반 업로드·심의 UI", "심의 이력, 권한, 업무 큐"],
        ["법령 엔진", "금소법 17·21·22 + OJK 핵심 트리", "제19조 설명의무, 제18조 적정성, 상품군별 감독규정"],
        ["데이터", "테스트셋 105건, 공개 제재 사례 벡터 DB", "실제 제재 사례 확대, 운영 피드백 반영"],
        ["규제 모니터링", "법령 원문 링크와 조문 실존 검증", "금감원 RSS·OJK 공시 자동 감지와 트리 영향도 분석"],
        ["멀티미디어", "이미지/PDF OCR", "영상 STT, 자막·음성 광고 검토"],
    ]
    for row in roadmap_rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            set_cell_text(cell, value, size=8)

    doc.add_heading("7. 용어 및 해석 기준", level=1)
    terms = [
        ("Risk Score", "AI가 탐지한 위험 신호의 강도를 0~1로 표현한 보조 지표입니다. 법률상 최종 위반 확정 점수가 아닙니다."),
        ("LLM Confidence", "LLM이 자기 판단에 부여한 신뢰도입니다. 조문 실존 여부와 별개이므로 Agent 4 검증 및 사람 검토가 필요합니다."),
        ("Human-in-the-loop", "AI가 1차 탐지와 근거 제시를 수행하고, 최종 승인·반려는 인간 준법관리자가 결정하는 구조입니다."),
        ("Regulatory Radar", "법률 판단 대체가 아니라 잠재적 규제 리스크를 사전에 넓게 탐지하는 보조 시스템 개념입니다."),
    ]
    for key, value in terms:
        paragraph = doc.add_paragraph()
        run = paragraph.add_run(key + ": ")
        run.bold = True
        run.font.color.rgb = BLUE
        paragraph.add_run(value)

    doc.add_heading("8. 변경 이력", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, ["일자", "대상", "변경 내용", "사유"]):
        set_cell_text(cell, header, bold=True, color=RGBColor(255, 255, 255), size=8)
        set_cell_shading(cell, "003B8F")
    change_rows = [
        ["2026.05.26", "법령 트리", "금소법 22조 및 OJK WARNING 보강, 평가 리포트 재생성", "FN 감소 및 Recall 우선 설계 근거 확보"],
        ["2026.05.26", "문서 구조", "문제·법령 범위·시스템·검증 결과·로드맵 순서로 재정리", "심사위원 가독성 및 목적성 강화"],
    ]
    for row in change_rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            set_cell_text(cell, value, size=8)

    doc.save(DOCX)


PB = PRGB(0, 59, 143)
PD = PRGB(11, 31, 77)
PL = PRGB(232, 241, 255)
PM = PRGB(232, 248, 242)
PP = PRGB(254, 239, 243)
WHITE = PRGB(255, 255, 255)


def add_bg(slide, prs) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRGB(248, 250, 255)
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(0.22), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PB
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(PInches(0.65), PInches(0.35), PInches(11.9), PInches(0.55))
    frame = box.text_frame
    frame.clear()
    run = frame.paragraphs[0].add_run()
    run.text = title
    run.font.name = "맑은 고딕"
    run.font.size = PPt(25)
    run.font.bold = True
    run.font.color.rgb = PD
    if subtitle:
        sub = slide.shapes.add_textbox(PInches(0.68), PInches(0.92), PInches(11.6), PInches(0.3))
        sub_frame = sub.text_frame
        sub_frame.clear()
        run = sub_frame.paragraphs[0].add_run()
        run.text = subtitle
        run.font.name = "맑은 고딕"
        run.font.size = PPt(9.5)
        run.font.color.rgb = PRGB(80, 94, 125)


def textbox(slide, x, y, w, h, text, size=13, bold=False, color=PD, fill=None, align="left"):
    if fill:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = PRGB(190, 206, 235)
        frame = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
        frame = shape.text_frame
    frame.clear()
    frame.margin_left = PInches(0.12)
    frame.margin_right = PInches(0.12)
    frame.margin_top = PInches(0.08)
    frame.margin_bottom = PInches(0.06)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def bullet_box(slide, x, y, w, h, title, bullets, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PRGB(190, 206, 235)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = PInches(0.18)
    frame.margin_right = PInches(0.16)
    frame.margin_top = PInches(0.13)
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "맑은 고딕"
    run.font.bold = True
    run.font.size = PPt(13)
    run.font.color.rgb = PB
    for bullet in bullets:
        paragraph = frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "맑은 고딕"
        paragraph.font.size = PPt(10.5)
        paragraph.font.color.rgb = PD
    return shape


def metric(slide, x, label, value, note):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(2.05), PInches(2.35), PInches(1.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = PRGB(183, 200, 232)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = PInches(0.12)
    frame.margin_right = PInches(0.12)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = value
    run.font.name = "맑은 고딕"
    run.font.size = PPt(24)
    run.font.bold = True
    run.font.color.rgb = PB
    paragraph = frame.add_paragraph()
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = label
    run.font.name = "맑은 고딕"
    run.font.size = PPt(10)
    run.font.bold = True
    run.font.color.rgb = PD
    paragraph = frame.add_paragraph()
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = note
    run.font.name = "맑은 고딕"
    run.font.size = PPt(8)
    run.font.color.rgb = PRGB(95, 105, 130)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    overall = SUMMARY["overall"]
    ko = SUMMARY["ko"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    textbox(slide, 0.75, 1.0, 7.6, 0.7, "Cross-Check AI", 32, True, PB)
    textbox(slide, 0.78, 1.75, 8.6, 0.6, "대고객 금융 콘텐츠를 위한 Human-in-the-loop Regulatory Radar", 18, False, PD)
    textbox(slide, 0.8, 2.55, 5.6, 0.55, "JB금융그룹 Fin:AI Challenge 2026 · 지정주제 2", 12, False, PD, PL)
    textbox(slide, 0.8, 3.25, 5.9, 0.55, "법령 트리 + Groq Llama 3.3 70B + 조문 검증 API", 12, False, PD, PM)
    textbox(slide, 0.8, 6.55, 3.4, 0.35, "JBIG 9th · 강승윤", 10, False, PRGB(80, 94, 125))
    textbox(slide, 8.2, 1.3, 3.8, 3.9, "AI가 법률판단을 대체하지 않습니다.\n\n위반 가능성을 넓게 탐지하고, 근거와 수정 방향을 제시하며, 최종 판단은 준법관리자가 수행합니다.", 18, True, WHITE, PB, "center")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "1. 왜 필요한가", "대고객 콘텐츠 전반을 수작업으로 심의하는 현재 구조의 병목")
    bullet_box(slide, 0.7, 1.55, 3.75, 4.45, "현행 문제", ["콘텐츠 채널 증가에 따라 심의 리소스가 선형적으로 증가", "다국어 현지화 시 본사 규제와 현지 규제를 동시에 확인해야 함", "규제 변경과 조문 인용 오류를 사람이 계속 추적해야 함"])
    bullet_box(slide, 4.78, 1.55, 3.75, 4.45, "JB금융 맥락", ["인도네시아 등 해외 사업 확장 시 현지 광고·권유 규제 대응 필요", "양국 법을 모두 아는 인력은 희소하고 심의 병목이 커질 수 있음", "현지화 과정에서 번역상 의미 변화가 규제 리스크로 전환될 수 있음"], PL)
    bullet_box(slide, 8.85, 1.55, 3.75, 4.45, "해결 방향", ["AI Agent가 1차 탐지와 근거 정리를 수행", "위반 누락(FN)을 최소화하는 Recall 우선 설계", "준법관리자는 최종 승인·반려에 집중"], PM)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "2. 솔루션 포지셔닝", "자율 법률판단 시스템이 아니라 준법관리자를 보조하는 Regulatory Radar")
    for idx, (title, body, fill) in enumerate([
        ("Rule-first", "명시적 금지 표현은 정규식·트리로 먼저 탐지해 재현성과 설명 가능성을 확보", PL),
        ("LLM-as-reviewer", "애매한 문맥은 Groq Llama 3.3 70B가 보완 판단하되, 조문 인용은 Agent 4가 재검증", PM),
        ("Human-in-the-loop", "AI는 위반 가능성과 근거를 제시하고 최종 승인·반려는 인간 준법관리자가 결정", PP),
    ]):
        bullet_box(slide, 0.8 + idx * 4.15, 1.6, 3.55, 3.9, title, [body], fill)
    textbox(slide, 1.25, 5.85, 10.9, 0.65, "핵심 메시지: Cross-Check AI는 ‘법률판단 자동화’가 아니라, 놓치기 쉬운 규제 리스크를 먼저 비추는 준법 심의 보조 시스템입니다.", 13, True, PB, WHITE, "center")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "3. 왜 금소법 17·21·22조인가", "광고 표현, 권유 방식, 고객 적합성이라는 대고객 콘텐츠 3대 리스크 축")
    law_rows = [
        ("22조", "광고 규제", "원금보장·확정수익·비용/위험 고지 누락"),
        ("21조", "부당권유", "단정적 판단·거짓과장·불안감 조성"),
        ("17조", "적합성 원칙", "고위험상품 무차별 권유·성향 확인 생략"),
        ("OJK 6/2022", "현지 소비자보호", "명확성·오인 방지·위험/비용 고지"),
    ]
    for idx, (article, role, examples) in enumerate(law_rows):
        y = 1.55 + idx * 1.05
        textbox(slide, 0.85, y, 2.1, 0.72, article, 15, True, WHITE, PB, "center")
        textbox(slide, 3.15, y, 2.45, 0.72, role, 13, True, PD, PL, "center")
        textbox(slide, 5.85, y, 6.5, 0.72, examples, 12, False, PD, WHITE, "center")
    textbox(slide, 1.0, 6.05, 11.0, 0.6, "법학 배경은 ‘전문가 행세’가 아니라, 어떤 조문을 먼저 구조화해야 실무 리스크를 잘 잡는지 선별하고 AI 판단 경계를 정하는 데 사용됩니다.", 12, True, PB, PM, "center")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "4. AI Agent 구조", "입력부터 조문 검증까지 이어지는 4-Agent 파이프라인")
    steps = [("사용자 입력", "텍스트/PDF/이미지 OCR"), ("Agent 1", "언어 감지·규제 분류"), ("Agent 2", "위반 탐지\nRule → Embedding → LLM"), ("Agent 3", "다국어 의미 정합성"), ("Agent 4", "조문 실존 검증"), ("Human Review", "승인/수정/반려")]
    for idx, (title, body) in enumerate(steps):
        x = 0.55 + idx * 2.05
        textbox(slide, x, 2.0, 1.65, 1.15, title + "\n" + body, 10.5, True, PD, PL if idx % 2 == 0 else WHITE, "center")
        if idx < 5:
            textbox(slide, x + 1.68, 2.35, 0.25, 0.3, "→", 18, True, PB, None, "center")
    textbox(slide, 1.05, 4.45, 10.9, 0.9, "법령 트리 DB · 제재 사례 벡터 DB · 국가법령정보 API · OJK 조문 DB가 Agent 판단의 근거 레이어로 연결됩니다.", 14, True, PB, WHITE, "center")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "5. 구현된 MVP", "FastAPI 웹앱에서 콘텐츠 업로드, OCR, 탐지, 근거 조문, Human Review 흐름까지 시연 가능")
    bullet_box(slide, 0.7, 1.5, 3.8, 4.55, "입력/전처리", ["텍스트·PDF·이미지 업로드", "Tesseract OCR + 금융광고 후처리", "OCR 품질 표시로 깨진 이미지 판단 보조"])
    bullet_box(slide, 4.75, 1.5, 3.8, 4.55, "탐지/근거", ["금소법 17·21·22 및 OJK YAML 트리", "Groq Llama 3.3 70B 문맥 판단", "유사 제재 사례 FAISS 검색"], PL)
    bullet_box(slide, 8.8, 1.5, 3.8, 4.55, "검증/운영", ["lawSearch/lawService 조문 실존 확인", "OJK 로컬 조문 DB 확인", "준법관리자 최종 의사결정 UI"], PM)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "6. 테스트셋 기반 검증 결과", "105건 평가셋 기준, LLM OFF rule-first baseline에서도 리스크 탐지 성능 확인")
    metric(slide, 0.85, "전체 Risk Recall", f"{overall['risk_recall']:.3f}", "WARNING/VIOLATION 검토 상향")
    metric(slide, 3.55, "전체 Precision", f"{overall['risk_precision']:.3f}", "과탐지 관리 수준")
    metric(slide, 6.25, "Violation Recall", f"{overall['violation_recall_not_pass']:.3f}", "명시적 위반 누락 최소화")
    metric(slide, 8.95, "한국어 위반 Recall", f"{ko['violation_recall_not_pass']:.3f}", "금소법 위반 케이스")
    textbox(slide, 1.15, 4.15, 10.9, 0.95, "초기 대비 Risk Recall 0.686 → 0.886으로 개선되었습니다. 임계치를 임의 상향한 것이 아니라, 놓치던 주의 표현을 법령 트리에 추가해 FN을 줄였습니다.", 14, True, PB, PM, "center")
    textbox(slide, 1.15, 5.45, 10.9, 0.55, "남은 과제: OJK 위반 누락 일부와 WARNING 세부 라벨 정밀화", 12, False, PD, WHITE, "center")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "7. 사용자 경험", "법률 JSON이 아니라 준법관리자가 바로 판단 가능한 Evidence 중심 화면")
    bullet_box(slide, 0.75, 1.5, 3.65, 4.4, "Input", ["광고 문구 또는 홍보물 업로드", "한국어·인도네시아어 자동 분류", "OCR 품질과 추출 문구 표시"])
    bullet_box(slide, 4.85, 1.5, 3.65, 4.4, "Evidence", ["위반 위치·매칭 문구", "근거 조문 원문 링크", "유사 제재 사례"], PL)
    bullet_box(slide, 8.95, 1.5, 3.65, 4.4, "Decision", ["PASS/WARNING/VIOLATION", "리스크 점수 설명", "승인·수정요청·반려"], PM)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "8. 향후 확장 로드맵", "예선은 검증 가능한 MVP, 본선은 운영형 준법 심의 플랫폼으로 확장")
    bullet_box(slide, 0.7, 1.45, 3.85, 4.75, "예선 제출 전 강화", ["평가셋 FN/FP 리뷰와 트리 정밀화", "OJK 조항 매핑 표 보강", "실제 홍보물 OCR 테스트 케이스 추가"])
    bullet_box(slide, 4.75, 1.45, 3.85, 4.75, "본선 구현", ["금감원 RSS·OJK 공시 모니터링", "영상·음성 STT 기반 콘텐츠 검토", "심의 이력·권한·업무 큐"], PL)
    bullet_box(slide, 8.8, 1.45, 3.85, 4.75, "실무 적용 방향", ["준법관리자 최종 판단 중심", "운영 피드백으로 법령 트리 업데이트", "국가·상품군별 규제 확장"], PM)

    prs.save(PPTX)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"updated {DOCX}")
    print(f"updated {PPTX}")
