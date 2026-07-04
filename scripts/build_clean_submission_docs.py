import json
from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt


ROOT = Path(__file__).resolve().parents[2]
PROJECT = ROOT / "cross_check_ai"
SPEC_OUT = ROOT / "CrossCheckAI_기능명세서.docx"
PPT_OUT = ROOT / "CrossCheckAI_MVP제안서.pptx"
SUMMARY = json.loads(
    (PROJECT / "data/evaluation/evaluation_report.summary.json").read_text(encoding="utf-8")
)


BLUE = RGBColor(0x00, 0x3B, 0x8F)
DEEP = RGBColor(0x0B, 0x1F, 0x4D)
LIGHT = "E8F1FF"
MINT = "E9F8F2"
BORDER = "B7C8E8"


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = tbl_pr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        tbl_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        elem = borders.find(qn("w:" + edge))
        if elem is None:
            elem = OxmlElement("w:" + edge)
            borders.append(elem)
        elem.set(qn("w:val"), "single")
        elem.set(qn("w:sz"), "6")
        elem.set(qn("w:space"), "0")
        elem.set(qn("w:color"), BORDER)


def set_cell(cell, text: str, bold=False, color=None, size=9, fill=None) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(text)
    r.bold = bold
    r.font.name = "맑은 고딕"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    r.font.size = Pt(size)
    if color:
        r.font.color.rgb = color
    if fill:
        set_cell_shading(cell, fill)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table(doc, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    for cell, header in zip(table.rows[0].cells, headers):
        set_cell(cell, header, True, RGBColor(255, 255, 255), 8.5, "003B8F")
    for row in rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            set_cell(cell, str(value), size=8)
    doc.add_paragraph()
    return table


def add_note(doc, title, lines, fill=LIGHT):
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_borders(table)
    cell = table.cell(0, 0)
    set_cell_shading(cell, fill)
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(title)
    r.bold = True
    r.font.name = "맑은 고딕"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    r.font.color.rgb = BLUE
    r.font.size = Pt(9.5)
    for line in lines:
        p = cell.add_paragraph()
        r = p.add_run("· " + line)
        r.font.name = "맑은 고딕"
        r._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
        r.font.size = Pt(9)
    doc.add_paragraph()


def build_docx():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.75)
    sec.right_margin = Inches(0.75)

    for style_name in ["Normal", "Heading 1", "Heading 2"]:
        style = doc.styles[style_name]
        style.font.name = "맑은 고딕"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    doc.styles["Normal"].font.size = Pt(9.5)
    doc.styles["Heading 1"].font.color.rgb = BLUE
    doc.styles["Heading 1"].font.size = Pt(14)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("JB금융그룹 Fin:AI Challenge 기능 명세서")
    r.bold = True
    r.font.name = "맑은 고딕"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    r.font.size = Pt(20)
    r.font.color.rgb = BLUE

    add_table(
        doc,
        ["팀명", "주제 구분", "팀원 정보(역할)", "작성일"],
        [["JBIG 9th", "지정주제 2 (준법자문가 AI Agent)", "강승윤 (PM · 기획 · 개발 · 디자인 총괄)", "2026.05.26"]],
    )

    doc.add_heading("1. 서비스 개요 (Service)", level=1)
    add_table(
        doc,
        ["항목", "내용"],
        [
            ["서비스명", "Cross-Check AI"],
            ["서비스 한줄 소개", "대고객 금융 콘텐츠의 규제 리스크를 법령 트리와 LLM으로 1차 탐지하고, 준법관리자가 최종 승인하는 Human-in-the-loop Regulatory Radar"],
            ["개발 목표", "수작업 중심 준법 심의의 지연·품질 편차·휴먼에러를 줄이고, 한국 금소법과 인도네시아 OJK 규정을 함께 검토할 수 있는 AI Agent 서비스 구현"],
            ["타겟 사용자", "금융사 준법관리자, 마케팅·영업·홍보 콘텐츠 제작 부서, 해외 법인 현지화 담당자"],
            ["기대 효과", "심의 리드타임 단축, 위반 누락(FN) 최소화, 조문 인용 오류 방지, 다국어 콘텐츠의 현지 규제 대응력 확보"],
        ],
    )

    doc.add_heading("2. 시스템 구성도 (Architecture)", level=1)
    add_note(
        doc,
        "사용자 입력 → AI Agent/모델 → 백엔드/API → DB/외부연동 → 프론트엔드/UI",
        [
            "사용자 입력: 텍스트, PDF, 이미지 OCR 기반 대고객 콘텐츠 업로드",
            "Agent 1: 언어 감지 및 금소법/OJK 규제 트리 라우팅",
            "Agent 2: Rule → Embedding → LLM 하이브리드 위반 탐지",
            "Agent 3: 한국어 원본과 인도네시아어 현지본의 의미·규제 정합성 비교",
            "Agent 4: 국가법령정보 API 및 OJK 조문 DB 기반 인용 조문 실존 검증",
            "FastAPI 웹앱: Evidence 중심 결과 화면과 준법관리자 승인/수정/반려 흐름 제공",
            "설계 근거: Norm AI식 규제-에이전트 접근, Stanford CodeX/John Nay의 Law Informs Code, 국내 법률 AI의 인용 검증 UX를 벤치마킹",
        ],
    )

    doc.add_heading("3. 핵심 기능 명세 (Feature Specification)", level=1)
    add_table(
        doc,
        ["기능명", "기능 설명", "입력/출력 데이터", "관련 기술 및 알고리즘", "구현 여부"],
        [
            ["콘텐츠 입력 및 규제 분류", "언어 감지 후 금소법/OJK 트리로 자동 라우팅", "텍스트/PDF/이미지 → 언어·규제 범위", "langdetect, OCR, FastAPI", "O"],
            ["위반 탐지 엔진", "법령을 YAML 의사결정 트리로 컴파일하고 명시적 위반은 Rule로 우선 탐지", "콘텐츠+트리 → PASS/WARNING/VIOLATION", "YAML DSL, Regex, Groq Llama 3.3 70B", "O"],
            ["근거 조문 검증", "인용 조문 실존 여부를 검증해 LLM 환각과 잘못된 인용 방지", "citation → 검증 결과·원문 링크", "국가법령정보 API, OJK 조문 DB", "O"],
            ["유사 제재 사례 검색", "위반 사유와 유사한 공개 제재 사례 검색", "판정 결과 → 유사 사례", "FAISS, sentence-transformers", "O"],
            ["다국어 정합성 검증", "한국어 원본과 인니어 현지본을 독립 심의 후 결과 차이 비교", "원본+현지본 → 규제 드리프트", "Termbase, 결과 diff", "O"],
            ["심의 결과 UI", "Evidence·Risk Score·수정 권고·Human Review 버튼 제공", "AI 판정 → 준법 심의 화면", "FastAPI 정적 웹 UI", "O"],
            ["규제 모니터링", "예선은 조문 검증, 본선은 RSS/OJK 공시 자동 모니터링으로 확장", "API/RSS → 개정 알림", "feedparser, scheduler", "△"],
        ],
    )

    doc.add_heading("4. 주요 기능 흐름도 (Flow)", level=1)
    add_note(
        doc,
        "핵심 이용 흐름",
        [
            "콘텐츠 업로드: 마케팅/영업 부서가 광고 문구, PDF, 이미지 홍보물을 업로드",
            "언어·규제 분류: Agent 1이 한국어/인도네시아어를 감지하고 금소법/OJK 트리를 선택",
            "위반 탐지: Agent 2가 명시적 표현은 Rule로, 애매한 문맥은 LLM으로 보완 검토",
            "근거 보강: 유사 제재 사례 검색과 Agent 4 조문 실존 검증으로 판단 근거 정리",
            "사람 검토: 준법관리자가 AI 결과를 확인하고 승인·수정 요청·반려 중 최종 결정",
        ],
        MINT,
    )

    doc.add_heading("5. 향후 발전 방향 (Future Work)", level=1)
    add_note(
        doc,
        "본선 및 실무 확장 계획",
        [
            "평가셋 기반 FN/FP 리뷰를 통해 금소법·OJK 트리 정밀화",
            "제19조 설명의무, 제18조 적정성 원칙, 상품군별 감독규정으로 법령 범위 확장",
            "금감원 RSS와 OJK 공시 자동 모니터링을 통한 개정 법령 영향도 분석",
            "영상·음성 광고의 STT 분석 기능은 본선 단계에서 구현",
            "운영 피드백을 법령 트리와 제재 사례 DB에 반영하는 지속 개선 구조 도입",
        ],
    )

    doc.add_heading("6. 부록 (Appendix)", level=1)
    add_note(
        doc,
        "참고자료 및 검증 결과",
        [
            "대회 상세주제 안내: 지정주제 2 준법자문가 AI Agent 서비스 개발",
            "법령 출처: 국가법령정보센터 Open API, POJK No. 6/POJK.07/2022",
            "해외 벤치마킹: Norm AI의 규제 준수 플랫폼은 법령·기관 맥락·규제 가이던스에 맞춘 reasoning alignment를 강조",
            "연구 근거: John J. Nay의 Law Informs Code는 법률 지식과 법적 추론을 AI 시스템에 구조적으로 내장하는 접근을 제시",
            "위험 근거: Stanford HAI/RegLab의 법률 LLM 환각 연구는 범용 LLM 단독 법률 판단의 위험성을 보여주며, Agent 4 조문 검증의 필요성을 뒷받침",
            "국내 벤치마킹: 슈퍼로이어·리걸테크 서비스처럼 판례·법령 인용에 하이퍼링크와 검증 장치를 제공하는 UX를 참고",
            "평가셋: 총 105건(한국어 60건, 금소법 17조 15건, OJK 30건)",
            f"검증 결과: Risk Recall {SUMMARY['overall']['risk_recall']:.3f}, Risk Precision {SUMMARY['overall']['risk_precision']:.3f}, Violation Recall {SUMMARY['overall']['violation_recall_not_pass']:.3f}",
            "Risk Score와 LLM Confidence는 최종 법률 판단이 아니라 준법관리자 검토 우선순위를 정하는 보조 지표",
        ],
        MINT,
    )

    doc.add_heading("7. 기능 변경이력 (Change Log)", level=1)
    add_table(
        doc,
        ["변경 일자", "변경 대상 기능", "변경 내용", "변경 사유"],
        [
            ["2026.05.26", "법령 트리", "금소법 22조 및 OJK WARNING 룰 보강", "Risk Recall 개선 및 FN 감소"],
            ["2026.05.26", "평가 리포트", "105건 테스트셋 기준 평가 결과 반영", "Recall 우선 설계 근거 확보"],
            ["2026.05.26", "제출 문서", "공식 7개 항목 양식 유지 + 겹침 없는 클린 레이아웃 재작성", "제출 형식 준수 및 가독성 개선"],
        ],
    )
    doc.save(SPEC_OUT)


P_BLUE = PRGB(0, 59, 143)
P_DEEP = PRGB(11, 31, 77)
P_GRAY = PRGB(82, 94, 120)
P_LIGHT = PRGB(232, 241, 255)
P_MINT = PRGB(232, 248, 242)
P_WHITE = PRGB(255, 255, 255)


def add_bg(slide, prs):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRGB(248, 250, 255)
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(0.22), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = P_BLUE
    bar.line.fill.background()
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0.22), 0, prs.slide_width - PInches(0.22), PInches(0.14))
    top.fill.solid()
    top.fill.fore_color.rgb = P_BLUE
    top.line.fill.background()


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(PInches(0.65), PInches(0.34), PInches(12.0), PInches(0.55))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.name = "맑은 고딕"
    r.font.size = PPt(23)
    r.font.bold = True
    r.font.color.rgb = P_DEEP
    if subtitle:
        box = slide.shapes.add_textbox(PInches(0.68), PInches(0.88), PInches(11.8), PInches(0.35))
        tf = box.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = subtitle
        r.font.name = "맑은 고딕"
        r.font.size = PPt(9.5)
        r.font.color.rgb = P_GRAY


def text_box(slide, x, y, w, h, text, size=12, bold=False, color=P_DEEP, fill=None, align="left"):
    if fill:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = PRGB(190, 206, 235)
    else:
        shape = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = PInches(0.13)
    tf.margin_right = PInches(0.13)
    tf.margin_top = PInches(0.08)
    tf.margin_bottom = PInches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "맑은 고딕"
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def card(slide, x, y, w, h, title, bullets, fill=P_WHITE):
    shape = text_box(slide, x, y, w, h, "", fill=fill)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = PInches(0.18)
    tf.margin_right = PInches(0.16)
    tf.margin_top = PInches(0.13)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "맑은 고딕"
    r.font.size = PPt(12.5)
    r.font.bold = True
    r.font.color.rgb = P_BLUE
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.name = "맑은 고딕"
        p.font.size = PPt(9.5)
        p.font.color.rgb = P_DEEP
    return shape


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    overall = SUMMARY["overall"]
    ko = SUMMARY["ko"]
    ojk = SUMMARY["ojk"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    text_box(slide, 0.75, 0.82, 7.8, 0.55, "JB금융그룹 Fin:AI Challenge", 20, True, P_BLUE)
    text_box(slide, 0.75, 1.45, 7.8, 0.7, "MVP 제안서", 30, True, P_DEEP)
    text_box(slide, 0.8, 2.35, 6.2, 0.55, "Cross-Check AI · JBIG 9th · 강승윤", 14, True, P_BLUE, P_LIGHT)
    text_box(slide, 0.8, 3.0, 8.5, 0.55, "대고객 금융 콘텐츠 준법 심의를 위한 Human-in-the-loop AI Agent", 13, False, P_DEEP)
    text_box(slide, 8.15, 1.35, 3.9, 3.7, "AI가 법률판단을 대체하지 않습니다.\n\n위반 가능성을 넓게 탐지하고, 근거와 수정 방향을 제시하며, 최종 판단은 준법관리자가 수행합니다.", 17, True, P_WHITE, P_BLUE, "center")

    slides = [
        (
            "예선 산출물 – 1. Summary",
            "JB금융그룹 Fin:AI Challenge",
            [
                ("문제", ["대고객 콘텐츠 준법 심의가 대부분 수작업", "다국어 콘텐츠 확장 시 심의 리소스가 선형 증가", "규제 변경·조문 인용 오류·휴먼에러 가능성 존재"]),
                ("해결", ["금소법·OJK 규정을 YAML 법령 트리로 구조화", "Rule-first + Groq Llama 3.3 70B + Agent 4 조문 검증", "AI는 1차 탐지, 최종 판단은 준법관리자"]),
                ("근거·검증", ["Norm AI·Stanford Law Informs Code 방향성 반영", f"105건 테스트셋 기준 Risk Recall {overall['risk_recall']:.3f}", f"Violation Recall {overall['violation_recall_not_pass']:.3f}, 한국어 위반 Recall {ko['violation_recall_not_pass']:.3f}"]),
            ],
            "핵심 메시지: Cross-Check AI는 법률판단을 대체하지 않고, 위반 가능성을 조기에 탐지하는 Regulatory Radar입니다.",
        ),
        (
            "2. 문제 정의(Problem Definition)",
            "예선 산출물 – 2. 문제 정의(Problem Definition)",
            [
                ("현행 분석", ["준법관리자가 광고·상품설명·SNS 등 대고객 콘텐츠를 수작업 검토", "외국어 콘텐츠도 동일 인력이 동일 방식으로 심의", "심의 지연으로 콘텐츠 배포 시점이 늦어질 수 있음"]),
                ("핵심 Pain Point", ["위반 누락(FN)은 실제 규제 위반 콘텐츠 배포로 이어질 수 있음", "규제 변경과 조문 인용 오류를 사람이 계속 추적해야 함", "다국어·다채널 확장 시 품질 편차와 휴먼에러 가능성 증가"]),
                ("해결 대상", ["1차 사용자: 금융사 준법관리자", "2차 사용자: 마케팅·영업·홍보 콘텐츠 제작 부서", "JB금융 인도네시아 등 해외 사업 확장에 따른 현지 규제 대응"]),
            ],
            None,
        ),
        (
            "3. 제안 솔루션 개요 (Solution Overview)",
            "예선 산출물 – 3. 제안 솔루션 개요 (Solution Overview)",
            [
                ("Agent 구조", ["Agent 1: 언어 감지·규제 분류", "Agent 2: Rule → Embedding → LLM 위반 탐지", "Agent 3: 다국어 의미·규제 정합성 비교", "Agent 4: 조문 실존 여부·원문 링크 검증"]),
                ("검증 가능한 법령 트리", ["Norm AI의 규제-에이전트 접근처럼 법령을 실행 가능한 판단 구조로 분해", "Stanford CodeX/John Nay의 Law Informs Code: 법률 지식과 추론을 AI에 구조화"]),
                ("Recall 우선 + 인용 검증", ["과탐지는 재검토 가능하지만 위반 누락은 과징금·제재·평판 손상으로 연결 가능", "국내 법률 AI처럼 조문·판례 인용은 링크와 검증 장치로 확인"]),
            ],
            "Cross-Check AI = 법령 트리 기반 준법 심의 Agent + Human-in-the-loop 검토 콘솔",
        ),
        (
            "4. 주요 기능 정의 (Key Features)",
            "예선 산출물 – 4. 주요 기능 정의 (Key Features)",
            [
                ("MVP 핵심 기능", ["콘텐츠 입력·분류: 텍스트/PDF/이미지 OCR, 한국어·인니어 자동 라우팅", "위반 탐지 엔진: 금소법 17·21·22 + OJK 트리", "근거 조문 검증: 국가법령정보 API와 OJK 조문 DB"]),
                ("보조 기능", ["유사 제재 검색: 공개 제재 사례 벡터 DB", "다국어 정합성: 원본·현지본 독립 심의 후 규제 드리프트 탐지", "Human Review UI: Risk Score 설명, 승인·수정요청·반려"]),
                ("구현 상태", ["FastAPI 웹앱 구현", "OCR·YAML 트리 엔진·Agent 4 조문 검증 구현", "105건 평가셋 리포트 반영"]),
            ],
            None,
        ),
        (
            "5. 데이터 및 기술 활용 (Data & Tech)",
            "예선 산출물 - 5. 데이터 및 기술 활용 (Data & Tech)",
            [
                ("활용 데이터", ["국가법령정보센터 Open API", "POJK No. 6/POJK.07/2022 조문 DB", "공개 제재 사례 DB + FAISS 벡터", "105건 테스트셋"]),
                ("벤치마킹·논문", ["Norm AI: regulated enterprise compliance workflow", "John J. Nay: Law Informs Code", "Stanford legal hallucination 연구: LLM 단독 판단 위험", "슈퍼로이어 등 국내 법률 AI: 인용 링크·검증 UX"]),
                ("검증 결과", [f"Risk Recall {overall['risk_recall']:.3f}", f"Risk Precision {overall['risk_precision']:.3f}", f"Violation Recall {overall['violation_recall_not_pass']:.3f}", f"OJK Risk Recall {ojk['risk_recall']:.3f}"]),
            ],
            None,
        ),
        (
            "6. 사용자 시나리오/유즈케이스 (User Scenario)",
            "예선 산출물 – 6. 사용자 시나리오/유즈케이스",
            [
                ("1. 콘텐츠 업로드", ["마케팅팀이 광고 문구·PDF·이미지 홍보물을 업로드", "OCR 품질과 추출 텍스트를 확인"]),
                ("2. AI 1차 심의", ["언어 감지 후 금소법/OJK 트리로 위반 가능성 탐지", "근거 조문·유사 제재 사례·매칭 문구 제시"]),
                ("3. 사람 최종 결정", ["준법관리자가 AI 결과를 확인", "승인·수정 요청·반려 중 최종 결정"]),
            ],
            "사용자 경험의 핵심은 법률 JSON이 아니라 준법관리자가 바로 판단 가능한 Evidence 중심 화면입니다.",
        ),
        (
            "7. 기대 효과 및 향후 확장성 (Expected Impact)",
            "예선 산출물 – 7. 기대 효과 및 향후 확장성",
            [
                ("예선 MVP 효과", ["심의 리스크를 사람이 보기 전에 1차 선별", "조문 인용 오류와 LLM 환각 방지", "테스트셋 기반 Recall 우선 설계 근거 확보"]),
                ("본선 확장", ["금감원 RSS·OJK 공시 자동 모니터링", "영상·음성 광고 STT 분석", "심의 이력, 권한, 업무 큐 기능"]),
                ("실무 적용성", ["준법관리자는 최종 판단에 집중", "운영 피드백으로 법령 트리 지속 개선", "국가·상품군별 규제 확장 가능"]),
            ],
            "Cross-Check AI는 JB금융의 동남아·다국어 콘텐츠 준법 심의를 수동 운영보다 빠르고 일관되게 보조합니다.",
        ),
    ]

    for idx, (title, subtitle, cards, footer) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, prs)
        add_title(slide, title, subtitle)
        for card_idx, (card_title, bullets) in enumerate(cards):
            x = 0.75 + card_idx * 4.05
            fill = [P_WHITE, P_LIGHT, P_MINT][card_idx % 3]
            card(slide, x, 1.55, 3.65, 4.25, card_title, bullets, fill)
        if footer:
            text_box(slide, 1.0, 6.1, 11.0, 0.55, footer, 12, True, P_BLUE, P_WHITE, "center")

    prs.save(PPT_OUT)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"saved {SPEC_OUT}")
    print(f"saved {PPT_OUT}")
