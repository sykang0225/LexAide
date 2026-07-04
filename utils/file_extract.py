"""
File text extraction for compliance review.

Priority:
1. Text-layer extraction for TXT/CSV/MD/PDF/PPTX/DOCX.
2. PDF text-layer quality is checked before compliance review.
3. If PDF text is unavailable/broken, OCR is used as a fallback input path.
4. OCR engines: EasyOCR first when available, Tesseract fallback.
"""
from __future__ import annotations

import importlib.util
import io
import logging
import os
import re
import subprocess
import sysconfig
import tempfile
import zipfile
from pathlib import Path
from shutil import which
from typing import Any
from xml.etree import ElementTree as ET

from utils.ocr_normalizer import normalize_ocr_text

logger = logging.getLogger(__name__)

_TESSDATA_DIR = Path(__file__).parent.parent / "data" / "tessdata"
_TESS_CANDIDATES = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]
_PADDLE_OCR: Any | None = None
_EASY_OCR: Any | None = None
_SURYA_MODEL: Any | None = None


class ExtractResult:
    def __init__(
        self,
        text: str,
        source_type: str,
        ocr_used: bool = False,
        note: str = "",
        ocr_quality: dict | None = None,
    ):
        self.text = text
        self.source_type = source_type
        self.ocr_used = ocr_used
        self.note = note
        self.ocr_quality = ocr_quality or {}

    def to_dict(self) -> dict:
        return {
            "text": self.text,
            "source_type": self.source_type,
            "ocr_used": self.ocr_used,
            "note": self.note,
            "ocr_quality": self.ocr_quality,
        }


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return data.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore").strip()


def _tesseract_path() -> str | None:
    for path in _TESS_CANDIDATES:
        if Path(path).exists():
            return path

    local_app_data = os.environ.get("LOCALAPPDATA", "")
    if local_app_data:
        candidate = Path(local_app_data) / "Programs" / "Tesseract-OCR" / "tesseract.exe"
        if candidate.exists():
            return str(candidate)

    from shutil import which

    return which("tesseract")


def tesseract_available() -> bool:
    return _tesseract_path() is not None


def paddle_available() -> bool:
    return bool(importlib.util.find_spec("paddleocr") and importlib.util.find_spec("paddle"))


def easyocr_available() -> bool:
    return bool(importlib.util.find_spec("easyocr") and importlib.util.find_spec("torch"))


def surya_available() -> bool:
    return bool(importlib.util.find_spec("surya"))


def ocr_available() -> bool:
    engine = os.environ.get("OCR_ENGINE", "auto").lower()
    if engine == "surya":
        return surya_available()
    if engine == "easyocr":
        return easyocr_available()
    if engine == "paddle":
        return paddle_available()
    if engine == "tesseract":
        return tesseract_available()
    return surya_available() or easyocr_available() or paddle_available() or tesseract_available()


def _clean_ocr_token(text: str) -> str:
    return (text or "").strip()


def _is_textish(text: str) -> bool:
    return bool(re.search(r"[가-힣A-Za-z0-9]", text or ""))


def _tesseract_ocr_data(data: bytes, min_conf: int = 45):
    import pytesseract
    from PIL import Image

    tesseract_path = _tesseract_path()
    if not tesseract_path:
        raise RuntimeError("TESSERACT_NOT_INSTALLED")
    pytesseract.pytesseract.tesseract_cmd = tesseract_path
    if _TESSDATA_DIR.exists():
        os.environ["TESSDATA_PREFIX"] = str(_TESSDATA_DIR)

    img = Image.open(io.BytesIO(data)).convert("RGB")

    # 속도 최적화: 과대 이미지는 OCR 전에 축소(긴 변 기준). 박스 좌표는 원본 스케일로 복원해
    # 부작위/현저성 분석(layout_risk.analyze_ocr_layout)의 위치·크기 좌표가 어긋나지 않게 한다.
    orig_w, orig_h = img.size
    max_side = int(os.environ.get("OCR_MAX_SIDE", "2200"))
    down = 1.0
    if max(orig_w, orig_h) > max_side:
        down = max_side / float(max(orig_w, orig_h))
        img = img.resize(
            (max(1, round(orig_w * down)), max(1, round(orig_h * down))),
            Image.LANCZOS,
        )
    inv = 1.0 / down

    result = pytesseract.image_to_data(
        img,
        lang=os.environ.get("TESS_LANGS", "kor+ind+eng"),
        config="--psm 6 --oem 1",
        output_type=pytesseract.Output.DICT,
    )

    boxes: list[dict] = []
    words: list[str] = []
    for i, raw in enumerate(result.get("text", [])):
        text = _clean_ocr_token(raw)
        try:
            conf = float(result["conf"][i])
        except (ValueError, TypeError):
            conf = -1

        if not text or conf < min_conf or not _is_textish(text):
            continue
        if len(text) == 1 and not re.search(r"[가-힣0-9]", text):
            continue

        boxes.append({
            "text": text,
            "x": int(round(int(result["left"][i]) * inv)),
            "y": int(round(int(result["top"][i]) * inv)),
            "w": int(round(int(result["width"][i]) * inv)),
            "h": int(round(int(result["height"][i]) * inv)),
            "conf": conf,
            "engine": "tesseract",
        })
        words.append(text)

    return normalize_ocr_text(" ".join(words)), boxes


def _surya_ocr_data(data: bytes):
    """Surya OCR — transformer 기반, 한국어·인도네시아어 고정밀."""
    global _SURYA_MODEL
    from PIL import Image
    from surya.ocr import run_ocr
    from surya.model.detection.model import load_model as load_det_model, load_processor as load_det_proc
    from surya.model.recognition.model import load_model as load_rec_model
    from surya.model.recognition.processor import load_processor as load_rec_proc

    if _SURYA_MODEL is None:
        _SURYA_MODEL = {
            "det_model": load_det_model(),
            "det_processor": load_det_proc(),
            "rec_model": load_rec_model(),
            "rec_processor": load_rec_proc(),
        }

    img = Image.open(io.BytesIO(data)).convert("RGB")
    langs = os.environ.get("SURYA_LANGS", "ko,id,en").split(",")
    results = run_ocr(
        [img],
        [langs],
        _SURYA_MODEL["det_model"],
        _SURYA_MODEL["det_processor"],
        _SURYA_MODEL["rec_model"],
        _SURYA_MODEL["rec_processor"],
    )

    boxes: list[dict] = []
    words: list[str] = []
    for line in (results[0].text_lines if results else []):
        text = _clean_ocr_token(line.text or "")
        if not text or not _is_textish(text):
            continue
        conf = round(float(line.confidence or 0) * 100, 1)
        bbox = line.bbox  # [x0, y0, x1, y1]
        boxes.append({
            "text": text,
            "x": int(bbox[0]), "y": int(bbox[1]),
            "w": max(1, int(bbox[2] - bbox[0])),
            "h": max(1, int(bbox[3] - bbox[1])),
            "conf": conf,
            "engine": "surya",
        })
        words.append(text)

    return normalize_ocr_text(" ".join(words)), boxes


def _get_easyocr_reader():
    global _EASY_OCR
    if _EASY_OCR is not None:
        return _EASY_OCR

    import easyocr

    # EasyOCR 제약: Korean은 English하고만 호환 (ko+id 조합 불가).
    # Indonesian 문서는 OCR_ENGINE=tesseract 또는 surya로 처리.
    langs = [
        lang.strip()
        for lang in os.environ.get("EASY_OCR_LANGS", "ko,en").split(",")
        if lang.strip()
    ]
    gpu = os.environ.get("EASY_OCR_GPU", "0") == "1"
    _EASY_OCR = easyocr.Reader(langs, gpu=gpu, verbose=False)
    return _EASY_OCR


def _easyocr_ocr_data(data: bytes, min_conf: int = 35):
    import numpy as np
    from PIL import Image

    reader = _get_easyocr_reader()
    img = Image.open(io.BytesIO(data)).convert("RGB")
    arr = np.array(img)
    results = reader.readtext(arr, detail=1, paragraph=False)

    boxes: list[dict] = []
    words: list[str] = []
    threshold = float(os.environ.get("EASY_OCR_MIN_CONF", max(0.20, min_conf / 100)))
    for item in results:
        if len(item) < 3:
            continue
        poly, raw_text, score = item[0], item[1], item[2]
        text = _clean_ocr_token(str(raw_text or ""))
        try:
            conf = float(score)
        except (TypeError, ValueError):
            conf = 0.0
        if not text or conf < threshold or not _is_textish(text):
            continue
        if not _is_polygon(poly):
            continue
        x, y, w, h = _bbox_from_polygon(poly)
        boxes.append({
            "text": text,
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "conf": round(conf * 100, 1),
            "engine": "easyocr",
        })
        words.append(text)

    return normalize_ocr_text(" ".join(words)), boxes


def _get_paddle_ocr():
    global _PADDLE_OCR
    if _PADDLE_OCR is not None:
        return _PADDLE_OCR

    from paddleocr import PaddleOCR

    lang = os.environ.get("PADDLE_OCR_LANG", "korean")
    init_attempts = (
        {"lang": lang, "use_angle_cls": True, "show_log": False},
        {"lang": lang, "use_angle_cls": True},
        {"lang": lang},
    )
    last_error: Exception | None = None
    for kwargs in init_attempts:
        try:
            _PADDLE_OCR = PaddleOCR(**kwargs)
            return _PADDLE_OCR
        except Exception as exc:  # PaddleOCR has changed constructor args across versions.
            last_error = exc
    raise RuntimeError(f"PADDLE_OCR_INIT_FAILED: {last_error}")


def _is_point(value: Any) -> bool:
    return (
        isinstance(value, (list, tuple))
        and len(value) >= 2
        and isinstance(value[0], (int, float))
        and isinstance(value[1], (int, float))
    )


def _is_polygon(value: Any) -> bool:
    return isinstance(value, (list, tuple)) and len(value) >= 4 and all(_is_point(p) for p in value[:4])


def _bbox_from_polygon(poly: Any) -> tuple[int, int, int, int]:
    xs = [float(p[0]) for p in poly]
    ys = [float(p[1]) for p in poly]
    left, top = min(xs), min(ys)
    right, bottom = max(xs), max(ys)
    return int(left), int(top), max(1, int(right - left)), max(1, int(bottom - top))


def _parse_paddle_result(result: Any) -> list[dict]:
    boxes: list[dict] = []

    def add(poly: Any, text: Any, score: Any) -> None:
        token = _clean_ocr_token(str(text or ""))
        if not token or not _is_textish(token):
            return
        try:
            conf = float(score) * 100 if float(score) <= 1 else float(score)
        except (TypeError, ValueError):
            conf = 0.0
        x, y, w, h = _bbox_from_polygon(poly)
        boxes.append({
            "text": token,
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "conf": conf,
            "engine": "paddleocr",
        })

    def visit(obj: Any) -> None:
        if obj is None:
            return

        if hasattr(obj, "json"):
            try:
                data = obj.json() if callable(obj.json) else obj.json
                visit(data)
                return
            except Exception:
                pass

        if isinstance(obj, dict):
            texts = obj.get("rec_texts") or obj.get("texts") or obj.get("text")
            scores = obj.get("rec_scores") or obj.get("scores") or obj.get("rec_score")
            polys = obj.get("rec_polys") or obj.get("dt_polys") or obj.get("boxes") or obj.get("polys")
            if isinstance(texts, list) and isinstance(polys, list):
                if not isinstance(scores, list):
                    scores = [0.0] * len(texts)
                for poly, text, score in zip(polys, texts, scores):
                    if _is_polygon(poly):
                        add(poly, text, score)
                return
            return

        if isinstance(obj, (list, tuple)):
            if len(obj) >= 2 and _is_polygon(obj[0]):
                payload = obj[1]
                if isinstance(payload, (list, tuple)) and len(payload) >= 2:
                    add(obj[0], payload[0], payload[1])
                    return
                if isinstance(payload, str):
                    add(obj[0], payload, 0.0)
                    return
            for item in obj:
                visit(item)

    visit(result)
    return boxes


def _paddle_ocr_data(data: bytes, min_conf: int = 45):
    import numpy as np
    from PIL import Image

    ocr = _get_paddle_ocr()
    img = Image.open(io.BytesIO(data)).convert("RGB")
    arr = np.array(img)

    if hasattr(ocr, "ocr"):
        try:
            result = ocr.ocr(arr, cls=True)
        except TypeError:
            result = ocr.ocr(arr)
    elif hasattr(ocr, "predict"):
        result = ocr.predict(arr)
    else:
        raise RuntimeError("PADDLE_OCR_METHOD_NOT_FOUND")

    boxes = [b for b in _parse_paddle_result(result) if float(b.get("conf", 0)) >= min_conf]
    text = normalize_ocr_text(" ".join(b["text"] for b in boxes))
    return text, boxes


def ocr_data(data: bytes, min_conf: int = 45):
    """Return (clean_text, boxes) from image bytes.

    우선순위 (auto): Surya → EasyOCR → PaddleOCR → Tesseract
    OCR_ENGINE 환경변수로 엔진 고정 가능.
    """
    engine = os.environ.get("OCR_ENGINE", "auto").lower()

    # ── Surya (최우선: transformer 기반, 한국어·인도네시아어 고정밀) ──
    if engine in ("auto", "surya") and surya_available():
        try:
            text, boxes = _surya_ocr_data(data)
            if text.strip() and boxes:
                return text, boxes
            logger.warning("[file_extract] Surya OCR returned empty; trying next engine.")
            if engine == "surya":
                return text, boxes
        except Exception as exc:
            logger.warning("[file_extract] Surya OCR failed: %s", exc)
            if engine == "surya":
                raise

    # ── EasyOCR ───────────────────────────────────────────────────────
    if engine in ("auto", "easyocr") and easyocr_available():
        try:
            text, boxes = _easyocr_ocr_data(data, min_conf=max(25, min_conf - 10))
            if text.strip() and boxes:
                return text, boxes
            logger.warning("[file_extract] EasyOCR returned empty/low-quality text; fallback will be tried.")
            if engine == "easyocr":
                return text, boxes
        except Exception as exc:
            logger.warning("[file_extract] EasyOCR failed: %s", exc)
            if engine == "easyocr":
                raise

    # ── Tesseract (auto에서 EasyOCR 저품질 시 우선 폴백 — 한·인니·영문 단일 패스) ──
    # PaddleOCR 모델이 손상된 환경이 많아, auto에서도 Tesseract를 Paddle보다 먼저 시도한다.
    if engine in ("auto", "tesseract") and tesseract_available():
        try:
            text, boxes = _tesseract_ocr_data(data, min_conf=min_conf)
            if (text.strip() and boxes) or engine == "tesseract":
                return text, boxes
            logger.warning("[file_extract] Tesseract returned empty; trying PaddleOCR.")
        except Exception as exc:
            logger.warning("[file_extract] Tesseract OCR failed: %s", exc)
            if engine == "tesseract":
                raise

    # ── PaddleOCR (최종 폴백) ─────────────────────────────────────────
    if engine in ("auto", "paddle") and paddle_available():
        try:
            return _paddle_ocr_data(data, min_conf=min_conf)
        except Exception as exc:
            logger.warning("[file_extract] PaddleOCR failed: %s", exc)
            if engine == "paddle":
                raise

    raise RuntimeError("OCR_ENGINE_UNAVAILABLE")


def assess_ocr_quality(data: bytes, text: str, boxes: list[dict]) -> dict:
    from PIL import Image

    try:
        img = Image.open(io.BytesIO(data))
        width, height = img.size
    except Exception:
        width, height = 0, 0

    confs = [float(b.get("conf", 0)) for b in boxes if b.get("conf") is not None]
    avg_conf = round(sum(confs) / len(confs), 1) if confs else 0.0
    word_count = len(boxes)
    text_len = len((text or "").strip())
    engine = boxes[0].get("engine", "unknown") if boxes else "unknown"

    issues = []
    if width and height and (width < 700 or height < 450):
        issues.append("이미지 해상도가 낮아 작은 글자 인식이 불안정할 수 있습니다.")
    if word_count < 8 or text_len < 40:
        issues.append("인식된 텍스트가 적어 전체 광고 내용을 대표하기 어렵습니다.")
    if avg_conf and avg_conf < 58:
        issues.append("OCR 평균 신뢰도가 낮습니다.")

    if issues:
        level = "LOW"
        message = "OCR 신뢰도가 낮습니다. 원본 PDF/PPTX의 텍스트 추출 또는 수동 확인이 필요합니다."
    elif avg_conf < 70:
        level = "MEDIUM"
        message = "OCR 결과를 사용할 수 있지만, 작은 글자나 주석은 수동 확인이 필요합니다."
    else:
        level = "HIGH"
        message = "OCR 결과가 심의 입력으로 사용할 수 있는 수준입니다."

    return {
        "level": level,
        "message": message,
        "issues": issues,
        "avg_conf": avg_conf,
        "word_count": word_count,
        "text_length": text_len,
        "width": width,
        "height": height,
        "engine": engine,
    }


def _pdf_text(data: bytes) -> str:
    import pdfplumber

    chunks = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            chunks.append(page.extract_text() or "")
    return normalize_ocr_text("\n".join(chunks).strip())


def _clean_pdf_text_artifacts(text: str) -> str:
    cleaned = text or ""
    cleaned = re.sub(r"\bn{2,}\b", " ", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


def _looks_broken_pdf_text(text: str) -> bool:
    """Detect PDF text layers whose font encoding extracts as repeated garbage."""
    cleaned = (text or "").strip()
    if len(cleaned) < 40:
        return False

    tokens = re.findall(r"[A-Za-z가-힣0-9]+", cleaned)
    if len(tokens) < 5:
        return True

    n_tokens = sum(1 for token in tokens if re.fullmatch(r"n{2,}", token.lower()))
    latin = re.findall(r"[A-Za-z]", cleaned)
    n_ratio = cleaned.lower().count("n") / max(1, len(latin))
    hangul_count = len(re.findall(r"[가-힣]", cleaned))

    # Some PDFs generated with embedded/subset fonts expose text as "nn nnnn ...".
    return n_tokens / len(tokens) >= 0.35 and n_ratio >= 0.45 and hangul_count < 5


def _looks_unusable_extracted_text(text: str) -> bool:
    cleaned = (text or "").strip()
    if len(cleaned) < 20:
        return True

    hangul_count = len(re.findall(r"[가-힣]", cleaned))
    tokens = re.findall(r"[A-Za-z가-힣0-9]+", cleaned)
    junk_tokens = {
        "BB", "BBB", "BBBB", "HB", "HHB", "MB", "EEE", "DD", "OO", "000", "0000", "00000"
    }
    meaningful_latin = [
        token for token in tokens
        if re.fullmatch(r"[A-Za-z]{3,}", token) and token.upper() not in junk_tokens
    ]
    # If Korean disappeared into tofu boxes, OCR often returns BB/HB/000-like
    # noise. Numbers alone are not enough for compliance review.
    return hangul_count < 8 and len(meaningful_latin) < 3


def _pdf_text_layer_quality(text: str) -> tuple[bool, str, str]:
    """Return (ok, cleaned_text, reason) for PDF text-layer quality."""
    raw = text or ""
    cleaned = _clean_pdf_text_artifacts(raw)
    stripped = cleaned.strip()

    if not stripped:
        return False, "", "empty_text_layer"
    if len(stripped) < 20:
        return False, stripped, "too_short_text_layer"

    if re.search(r"[□■�]{2,}", raw):
        return False, stripped, "replacement_character_detected"
    if _looks_broken_pdf_text(raw):
        return False, stripped, "broken_font_encoding"

    # ★ 텍스트레이어에 실질 콘텐츠가 충분하면 즉시 신뢰 → OCR 절대 미실행 (정상 텍스트 글자 깨짐 방지)
    _hangul = len(re.findall(r"[가-힣]", stripped))
    _latin = len([t for t in re.findall(r"[A-Za-zÀ-ÿ]+", stripped) if len(t) >= 3])
    if _hangul >= 15 or _latin >= 8:
        return True, stripped, ""

    total_chars = len(re.sub(r"\s+", "", stripped))
    normal_chars = len(re.findall(r"[가-힣A-Za-zÀ-ÿ0-9%₩$.,:;()!?+\-/]", stripped))
    normal_ratio = normal_chars / max(1, total_chars)
    if normal_ratio < 0.65:
        return False, stripped, "low_normal_character_ratio"

    tokens = re.findall(r"[A-Za-zÀ-ÿ가-힣0-9]+", stripped)
    junk_tokens = {"BB", "BBB", "BBBB", "HB", "HHB", "MB", "EEE", "DD", "OO"}
    hangul_count = len(re.findall(r"[가-힣]", stripped))
    meaningful_latin = [
        token for token in tokens
        if re.fullmatch(r"[A-Za-zÀ-ÿ]{3,}", token) and token.upper() not in junk_tokens
    ]

    # Korean documents need Korean text. Indonesian/English documents can pass
    # with meaningful Latin words. Pure numbers and placeholder glyphs cannot.
    if hangul_count < 3 and len(meaningful_latin) < 3:
        return False, stripped, "insufficient_meaningful_language_text"

    return True, stripped, ""


def _pdf_text_layer_error_note() -> str:
    return (
        "업로드한 PDF의 텍스트 레이어를 정상적으로 읽을 수 없습니다. "
        "정확한 준법 검증을 위해 텍스트 선택이 가능한 PDF, 원본 문서, "
        "또는 광고 문구 텍스트를 함께 제출해주세요."
    )


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []

    def read_shape(shape) -> None:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                read_shape(child)
            return
        if getattr(shape, "has_text_frame", False):
            text = (shape.text or "").strip()
            if text:
                chunks.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = (cell.text or "").strip()
                    if text:
                        chunks.append(text)

    for i, slide in enumerate(prs.slides, start=1):
        slide_chunks_before = len(chunks)
        for shape in slide.shapes:
            read_shape(shape)
        if len(chunks) > slide_chunks_before:
            chunks.insert(slide_chunks_before, f"[Slide {i}]")

    return normalize_ocr_text("\n".join(chunks).strip())


def _docx_text(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    chunks: list[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            chunks.append(text)
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                chunks.append(" | ".join(row_text))
    return normalize_ocr_text("\n".join(chunks).strip())


def _hwp5txt_command() -> str | None:
    cmd = which("hwp5txt")
    if cmd:
        return cmd

    exe_name = "hwp5txt.exe" if os.name == "nt" else "hwp5txt"
    candidates: list[Path] = []

    scripts_dir = sysconfig.get_path("scripts")
    if scripts_dir:
        candidates.append(Path(scripts_dir) / exe_name)

    user_base = sysconfig.get_config_var("userbase")
    if user_base:
        candidates.append(Path(user_base) / ("Scripts" if os.name == "nt" else "bin") / exe_name)

    appdata = os.environ.get("APPDATA")
    if appdata:
        candidates.append(Path(appdata) / "Python" / "Python312" / "Scripts" / exe_name)

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def hwp_available() -> bool:
    return _hwp5txt_command() is not None


def _hwp_text(data: bytes) -> str:
    cmd = _hwp5txt_command()
    if not cmd:
        raise RuntimeError("HWP5TXT_NOT_INSTALLED")

    temp_path: Path | None = None
    output_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".hwp", delete=False) as tmp:
            tmp.write(data)
            temp_path = Path(tmp.name)

        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as out:
            output_path = Path(out.name)

        completed = subprocess.run(
            [cmd, "--output", str(output_path), str(temp_path)],
            capture_output=True,
            timeout=30,
            check=False,
        )
        if completed.returncode != 0:
            err = _decode_text(completed.stderr or b"")
            raise RuntimeError(f"HWP5TXT_FAILED: {err or completed.returncode}")

        if output_path.exists():
            text = _decode_text(output_path.read_bytes())
        else:
            text = _decode_text(completed.stdout or b"")
        return normalize_ocr_text(text)
    finally:
        for path in (temp_path, output_path):
            if path:
                try:
                    path.unlink(missing_ok=True)
                except Exception:
                    pass


def _hwpx_text(data: bytes) -> str:
    chunks: list[str] = []
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = sorted(n for n in zf.namelist() if n.lower().endswith(".xml"))
        section_names = [
            n for n in names
            if re.search(r"(^|/)contents/section\d+\.xml$", n.lower())
        ]
        target_names = section_names or names

        for name in target_names:
            try:
                root = ET.fromstring(zf.read(name))
            except Exception:
                continue
            for value in root.itertext():
                text = (value or "").strip()
                if text:
                    chunks.append(text)

    return normalize_ocr_text(" ".join(chunks).strip())


def _ocr_image(data: bytes) -> tuple[str, list[dict]]:
    text, boxes = ocr_data(data)
    return text.strip(), boxes


def _pdf_pages_to_images(data: bytes):
    try:
        import fitz

        doc = fitz.open(stream=data, filetype="pdf")
        zoom = float(os.environ.get("PDF_OCR_ZOOM", "3.0"))
        matrix = fitz.Matrix(zoom, zoom)
        for page in doc:
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            yield pix.tobytes("png")
        return
    except Exception as exc:
        logger.warning("[file_extract] PyMuPDF PDF rendering failed: %s", exc)

    import pdf2image

    for page in pdf2image.convert_from_bytes(data):
        yield _pil_to_bytes(page)


def _ocr_pdf_data(data: bytes):
    chunks = []
    all_boxes: list[dict] = []
    page_no = 0
    for page_bytes in _pdf_pages_to_images(data):
        page_no += 1
        text, boxes = _ocr_image(page_bytes)
        if text:
            chunks.append(text)
        for box in boxes:
            item = dict(box)
            item["page"] = page_no
            all_boxes.append(item)
    return normalize_ocr_text("\n".join(chunks).strip()), all_boxes


def _ocr_scanned_pdf(data: bytes) -> str:
    text, _ = _ocr_pdf_data(data)
    return text


def _ocr_text_is_usable(text: str, boxes: list[dict]) -> tuple[bool, str]:
    if _looks_unusable_extracted_text(text):
        return False, "ocr_unusable_text"
    if len((text or "").strip()) < 40:
        return False, "ocr_text_too_short"
    if len(boxes) < 8:
        return False, "ocr_too_few_tokens"
    confs = [float(b.get("conf", 0)) for b in boxes if b.get("conf") is not None]
    avg_conf = sum(confs) / len(confs) if confs else 0.0
    if avg_conf < float(os.environ.get("PDF_OCR_MIN_AVG_CONF", "35")):
        return False, "ocr_low_confidence"
    return True, ""


def extract_text(data: bytes, filename: str) -> ExtractResult:
    ext = Path(filename or "").suffix.lower()

    if ext in (".txt", ".md", ".csv"):
        return ExtractResult(_decode_text(data), "txt")

    if ext == ".docx":
        try:
            text = _docx_text(data)
            return ExtractResult(text, "docx", note="DOCX 텍스트 레이어에서 추출했습니다.")
        except Exception as exc:
            logger.warning("[file_extract] DOCX extraction failed: %s", exc)
            return ExtractResult("", "docx", note="DOCX 텍스트 추출 중 오류가 발생했습니다.")

    if ext == ".pptx":
        try:
            text = _pptx_text(data)
            return ExtractResult(text, "pptx", note="PPTX 텍스트 레이어에서 추출했습니다.")
        except Exception as exc:
            logger.warning("[file_extract] PPTX extraction failed: %s", exc)
            return ExtractResult("", "pptx", note="PPTX 텍스트 추출 중 오류가 발생했습니다.")

    if ext == ".hwp":
        try:
            text = _hwp_text(data)
            if text:
                return ExtractResult(text, "hwp", note="HWP 문서에서 hwp5txt로 텍스트를 추출했습니다.")
            return ExtractResult("", "hwp", note="HWP 문서에서 심의 가능한 텍스트를 찾지 못했습니다.")
        except RuntimeError as exc:
            logger.warning("[file_extract] HWP extraction failed: %s", exc)
            if "HWP5TXT_NOT_INSTALLED" in str(exc):
                return ExtractResult(
                    "",
                    "hwp",
                    note="HWP 추출 도구(hwp5txt)가 설치되어 있지 않습니다. requirements 설치 후 다시 시도하세요.",
                )
            return ExtractResult("", "hwp", note="HWP 텍스트 추출 중 오류가 발생했습니다.")
        except Exception as exc:
            logger.warning("[file_extract] HWP extraction failed: %s", exc)
            return ExtractResult("", "hwp", note="HWP 텍스트 추출 중 오류가 발생했습니다.")

    if ext == ".hwpx":
        try:
            text = _hwpx_text(data)
            if text:
                return ExtractResult(text, "hwpx", note="HWPX 문서 XML에서 텍스트를 추출했습니다.")
            return ExtractResult("", "hwpx", note="HWPX 문서에서 심의 가능한 텍스트를 찾지 못했습니다.")
        except Exception as exc:
            logger.warning("[file_extract] HWPX extraction failed: %s", exc)
            return ExtractResult("", "hwpx", note="HWPX 텍스트 추출 중 오류가 발생했습니다.")

    if ext == ".pdf":
        quality_reason = "unknown"
        raw_text = ""
        try:
            raw_text = _pdf_text(data)
            ok, text, quality_reason = _pdf_text_layer_quality(raw_text)
            if ok:
                return ExtractResult(text, "pdf_text", note="PDF 텍스트 레이어에서 추출했습니다.")
            logger.info("[file_extract] PDF text layer rejected: %s", quality_reason)
        except Exception as exc:
            quality_reason = "text_layer_extraction_failed"
            logger.warning("[file_extract] PDF text extraction failed: %s", exc)

        if ocr_available():
            try:
                text, boxes = _ocr_pdf_data(data)
                usable, ocr_reason = _ocr_text_is_usable(text, boxes)
                if usable:
                    confs = [float(b.get("conf", 0)) for b in boxes if b.get("conf") is not None]
                    avg_conf = round(sum(confs) / len(confs), 1) if confs else 0.0
                    return ExtractResult(
                        text,
                        "pdf_ocr",
                        ocr_used=True,
                        note="PDF 텍스트 레이어 품질이 낮아 OCR로 보조 추출했습니다. 최종 제출 전 원문 문구 확인을 권장합니다.",
                        ocr_quality={
                            "level": "MEDIUM" if avg_conf < 70 else "HIGH",
                            "message": "텍스트 레이어 불량 PDF를 OCR로 보조 추출했습니다.",
                            "issues": [f"text_layer:{quality_reason}"],
                            "avg_conf": avg_conf,
                            "word_count": len(boxes),
                            "text_length": len(text or ""),
                            "engine": boxes[0].get("engine", "ocr") if boxes else "ocr",
                        },
                    )
                quality_reason = f"{quality_reason};{ocr_reason}"
            except Exception as exc:
                logger.warning("[file_extract] PDF OCR fallback failed: %s", exc)
                quality_reason = f"{quality_reason};ocr_failed"

        return ExtractResult(
            "",
            "pdf_text_bad_quality",
            ocr_used=False,
            note=_pdf_text_layer_error_note(),
            ocr_quality={
                "level": "LOW",
                "message": _pdf_text_layer_error_note(),
                "issues": [quality_reason],
                "engine": "pdf_text_layer",
                "text_length": len(raw_text or ""),
            },
        )

    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
        try:
            text, boxes = _ocr_image(data)
            quality = assess_ocr_quality(data, text, boxes)
            if not text:
                return ExtractResult(
                    "",
                    "image",
                    ocr_used=True,
                    note="이미지에서 글자를 인식하지 못했습니다. 원본 PDF/PPTX 또는 고해상도 이미지를 사용하세요.",
                    ocr_quality=quality,
                )
            return ExtractResult(text, "image", ocr_used=True, note=quality["message"], ocr_quality=quality)
        except Exception as exc:
            logger.warning("[file_extract] image OCR failed: %s", exc)
            return ExtractResult(
                "",
                "image",
                ocr_used=False,
                note="이미지 OCR 엔진을 찾지 못했거나 처리 중 오류가 발생했습니다. Tesseract 또는 PaddleOCR 설치가 필요합니다.",
            )

    return ExtractResult(
        "",
        ext.lstrip(".") or "unknown",
        note=f"지원하지 않는 파일 형식입니다: {ext}. 지원 형식: txt, md, csv, pdf, pptx, docx, hwp, hwpx, png, jpg",
    )


def _pil_to_bytes(image) -> bytes:
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    return buf.getvalue()
